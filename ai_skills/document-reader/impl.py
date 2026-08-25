import json
import os

from core.filesystem_ops import _build_error, _build_ok, record_full_read_state, resolve_path


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".xls", ".pdf"}


def _bundled_dependency_error(package_name, exc):
    raise RuntimeError(
        f"随应用安装的文档读取组件缺少 {package_name}。请重新安装完整的 Cowork 分发包。"
    ) from exc


def _get_openpyxl():
    try:
        import openpyxl

        return openpyxl
    except ImportError as exc:
        _bundled_dependency_error("openpyxl", exc)


def _get_docx_document():
    try:
        from docx import Document

        return Document
    except ImportError as exc:
        _bundled_dependency_error("python-docx", exc)


def _get_presentation():
    try:
        from pptx import Presentation

        return Presentation
    except ImportError as exc:
        _bundled_dependency_error("python-pptx", exc)


def _get_pdf_reader():
    try:
        from pypdf import PdfReader

        return PdfReader
    except ImportError as exc:
        _bundled_dependency_error("pypdf", exc)


def _truncate_text(text, max_chars):
    text = str(text or "")
    try:
        limit = int(max_chars or 60000)
    except Exception:
        limit = 60000
    if limit <= 0 or len(text) <= limit:
        return text, False
    return text[:limit], True


def _parse_pages(pages, total_pages):
    if pages is None:
        return list(range(total_pages)), None
    text = str(pages).strip()
    if not text:
        return list(range(total_pages)), None
    if "-" in text:
        start_raw, end_raw = text.split("-", 1)
        try:
            start_page = int(start_raw)
            end_page = int(end_raw)
        except Exception:
            return None, "pages must be in 'N' or 'N-M' format."
        if start_page <= 0 or end_page <= 0 or start_page > end_page:
            return None, "Invalid page range."
        first = start_page - 1
        last = min(end_page, total_pages)
        if first >= total_pages:
            return None, "Page range exceeds PDF page count."
        return list(range(first, last)), None
    try:
        single = int(text)
    except Exception:
        return None, "pages must be in 'N' or 'N-M' format."
    if single <= 0 or single > total_pages:
        return None, "Requested page is out of range."
    return [single - 1], None


def _read_docx(abs_path):
    Document = _get_docx_document()
    doc = Document(abs_path)
    sections = []
    for index, paragraph in enumerate(doc.paragraphs, start=1):
        if paragraph.text:
            sections.append({"type": "paragraph", "index": index, "text": paragraph.text})
    content = "\n".join(item["text"] for item in sections)
    return content, sections, {"paragraph_count": len(doc.paragraphs)}


def _read_pptx(abs_path):
    Presentation = _get_presentation()
    presentation = Presentation(abs_path)
    sections = []
    chunks = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slide_text = "\n".join(texts)
        sections.append({"type": "slide", "index": index, "text": slide_text})
        chunks.append(f"Slide {index}:\n{slide_text}")
    return "\n\n".join(chunks), sections, {"slide_count": len(presentation.slides)}


def _read_xlsx(abs_path, sheet_name=None):
    openpyxl = _get_openpyxl()
    workbook = openpyxl.load_workbook(abs_path, data_only=True)
    sheet_names = list(workbook.sheetnames)
    selected_names = [sheet_name] if sheet_name else sheet_names
    sections = []
    chunks = []
    for name in selected_names:
        if name not in workbook.sheetnames:
            raise ValueError(f"Sheet '{name}' not found.")
        worksheet = workbook[name]
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            rows.append("\t".join("" if cell is None else str(cell) for cell in row))
        text = "\n".join(rows)
        sections.append({"type": "sheet", "name": name, "text": text})
        chunks.append(f"Sheet {name}:\n{text}")
    return "\n\n".join(chunks), sections, {"sheet_names": sheet_names}


def _read_pdf(abs_path, pages=None):
    PdfReader = _get_pdf_reader()
    reader = PdfReader(abs_path)
    page_indices, error = _parse_pages(pages, len(reader.pages))
    if error:
        raise ValueError(error)
    sections = []
    chunks = []
    for page_index in page_indices:
        text = reader.pages[page_index].extract_text() or ""
        sections.append({"type": "page", "index": page_index + 1, "text": text})
        chunks.append(f"Page {page_index + 1}:\n{text}")
    return "\n\n".join(chunks), sections, {"total_pages": len(reader.pages), "returned_pages": len(page_indices)}


def document_read(workspace_dir, path, sheet_name=None, pages=None, max_chars=60000, _context=None):
    """Read DOCX, PPTX, XLSX, XLS, or PDF through one unified structured-document tool."""
    action = "document_read"
    ext = os.path.splitext(str(path or ""))[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return _build_error(
            action,
            "unsupported_format",
            "document_read supports DOCX, PPTX, XLSX, XLS, and PDF only. Use text_file_read for plain text files.",
            path=path,
        )
    abs_path, rel_path, error = resolve_path(workspace_dir, path, context=_context, action=action, must_exist=True)
    if error:
        return error
    if not os.path.isfile(abs_path):
        return _build_error(action, "not_a_file", "Path is not a file.", path=rel_path)

    try:
        if ext == ".docx":
            content, sections, metadata = _read_docx(abs_path)
        elif ext == ".pptx":
            content, sections, metadata = _read_pptx(abs_path)
        elif ext in {".xlsx", ".xls"}:
            content, sections, metadata = _read_xlsx(abs_path, sheet_name=sheet_name)
        else:
            content, sections, metadata = _read_pdf(abs_path, pages=pages)
        content, truncated = _truncate_text(content, max_chars)
        record_full_read_state(abs_path, _context)
        return _build_ok(
            action,
            {
                "path": rel_path,
                "format": ext.lstrip("."),
                "content": content,
                "metadata": metadata,
                "truncated": truncated,
                "sections": sections,
            },
        )
    except Exception as exc:
        return _build_error(action, "read_failed", str(exc), path=rel_path)


TOOL_EXPORTS = [
    {
        "name": "document_read",
        "handler": document_read,
        "description": "Read DOCX, PPTX, XLSX, XLS, or PDF files. Use text_file_read for plain text files.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Workspace-relative structured document path."},
                "sheet_name": {"type": "string", "description": "Optional XLSX/XLS sheet name."},
                "pages": {"type": "string", "description": "Optional PDF page selector, such as '1' or '2-4'."},
                "max_chars": {"type": "integer", "description": "Maximum content characters to return."},
            },
            "required": ["path"],
        },
        "read_only": True,
        "search_hint": "read docx pptx xlsx xls pdf office document structured",
    }
]
