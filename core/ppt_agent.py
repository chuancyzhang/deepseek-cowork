import os
import posixpath
import zipfile
from dataclasses import dataclass
from urllib.parse import unquote, urlparse
from xml.etree import ElementTree


PPT_AGENT_STRATEGY_AUTO = "auto"
PPT_AGENT_STRATEGY_DEFAULT = "default"
PPT_AGENT_STRATEGY_GUIZANG = "guizang"
PPT_AGENT_STRATEGY_FRONTEND_SLIDES = "frontend_slides"
PPT_AGENT_STRATEGY_HUASHU = "huashu"

PPT_AGENT_STRATEGY_CHOICES = {
    PPT_AGENT_STRATEGY_AUTO,
    PPT_AGENT_STRATEGY_DEFAULT,
    PPT_AGENT_STRATEGY_GUIZANG,
    PPT_AGENT_STRATEGY_FRONTEND_SLIDES,
    PPT_AGENT_STRATEGY_HUASHU,
}

PPT_AGENT_PREFERENCE_AUTO = "auto"
PPT_AGENT_PREFERENCE_WEB = "web"
PPT_AGENT_PREFERENCE_TECH = "tech"
PPT_AGENT_PREFERENCE_BUSINESS = "business"
PPT_AGENT_PREFERENCE_TEMPLATE = "template"

PPT_AGENT_PREFERENCES = {
    PPT_AGENT_PREFERENCE_AUTO,
    PPT_AGENT_PREFERENCE_WEB,
    PPT_AGENT_PREFERENCE_TECH,
    PPT_AGENT_PREFERENCE_BUSINESS,
    PPT_AGENT_PREFERENCE_TEMPLATE,
}

PPT_AGENT_OUTPUT_HTML = "html"
PPT_AGENT_OUTPUT_PPTX = "pptx"
PPT_AGENT_OUTPUT_FORMATS = {
    PPT_AGENT_OUTPUT_HTML,
    PPT_AGENT_OUTPUT_PPTX,
}


@dataclass(frozen=True)
class PptHtmlCapability:
    key: str
    name: str
    skill_name: str
    summary: str
    suitable_for: str
    output: str = "HTML"


PPT_HTML_CAPABILITIES = {
    PPT_AGENT_STRATEGY_GUIZANG: PptHtmlCapability(
        key=PPT_AGENT_STRATEGY_GUIZANG,
        name="Guizang PPT Skill",
        skill_name="guizang-ppt-skill",
        summary="HTML 横向翻页 PPT / PPT 配图 / 强视觉表达",
        suitable_for="网页演示、视觉分享、封面图、风格化内容",
    ),
    PPT_AGENT_STRATEGY_FRONTEND_SLIDES: PptHtmlCapability(
        key=PPT_AGENT_STRATEGY_FRONTEND_SLIDES,
        name="Frontend Slides",
        skill_name="frontend-slides",
        summary="前端技术生成 HTML Slides / 适合产品和技术演示",
        suitable_for="产品发布、技术分享、网页化演示、交互式演示",
    ),
    PPT_AGENT_STRATEGY_HUASHU: PptHtmlCapability(
        key=PPT_AGENT_STRATEGY_HUASHU,
        name="Huashu Design",
        skill_name="huashu-design",
        summary="高审美设计型 HTML PPT / 商业视觉表达",
        suitable_for="商业汇报、路演、品牌提案、发布会、产品介绍",
    ),
}

_GUIZANG_KEYWORDS = (
    "网页ppt",
    "网页 ppt",
    "网页演示",
    "横向翻页",
    "封面图",
    "配图",
    "视觉分享",
    "内容分享",
    "强风格",
    "风格化",
)
_FRONTEND_SLIDES_KEYWORDS = (
    "技术分享",
    "产品发布",
    "交互式",
    "前端",
    "slides",
    "slide",
    "网页化",
    "发布演示",
)
_HUASHU_KEYWORDS = (
    "高级感",
    "高审美",
    "设计感",
    "商业汇报",
    "路演",
    "品牌提案",
    "发布会",
    "产品介绍",
    "咨询风",
    "好看",
)


def normalize_ppt_agent_strategy(value):
    text = str(value or "").strip().lower().replace("-", "_")
    aliases = {
        "guizang_ppt_skill": PPT_AGENT_STRATEGY_GUIZANG,
        "frontend": PPT_AGENT_STRATEGY_FRONTEND_SLIDES,
        "frontend_slides": PPT_AGENT_STRATEGY_FRONTEND_SLIDES,
        "huashu_design": PPT_AGENT_STRATEGY_HUASHU,
        "ppt_agent": PPT_AGENT_STRATEGY_DEFAULT,
    }
    text = aliases.get(text, text)
    return text if text in PPT_AGENT_STRATEGY_CHOICES else PPT_AGENT_STRATEGY_AUTO


def normalize_ppt_agent_preference(value):
    text = str(value or "").strip().lower().replace("-", "_")
    return text if text in PPT_AGENT_PREFERENCES else PPT_AGENT_PREFERENCE_AUTO


def normalize_ppt_agent_output_format(value):
    text = str(value or "").strip().lower()
    return text if text in PPT_AGENT_OUTPUT_FORMATS else PPT_AGENT_OUTPUT_PPTX


def choose_ppt_agent_strategy(request_text, preference=PPT_AGENT_PREFERENCE_AUTO, explicit_strategy=PPT_AGENT_STRATEGY_AUTO, template_file=""):
    explicit = normalize_ppt_agent_strategy(explicit_strategy)
    if explicit != PPT_AGENT_STRATEGY_AUTO:
        return explicit
    preference = normalize_ppt_agent_preference(preference)
    if preference == PPT_AGENT_PREFERENCE_WEB:
        return PPT_AGENT_STRATEGY_GUIZANG
    if preference == PPT_AGENT_PREFERENCE_TECH:
        return PPT_AGENT_STRATEGY_FRONTEND_SLIDES
    if preference == PPT_AGENT_PREFERENCE_BUSINESS:
        return PPT_AGENT_STRATEGY_HUASHU
    if preference == PPT_AGENT_PREFERENCE_TEMPLATE or template_file:
        return PPT_AGENT_STRATEGY_DEFAULT
    text = str(request_text or "").strip().lower()
    if any(keyword in text for keyword in _HUASHU_KEYWORDS):
        return PPT_AGENT_STRATEGY_HUASHU
    if any(keyword in text for keyword in _FRONTEND_SLIDES_KEYWORDS):
        return PPT_AGENT_STRATEGY_FRONTEND_SLIDES
    if any(keyword in text for keyword in _GUIZANG_KEYWORDS):
        return PPT_AGENT_STRATEGY_GUIZANG
    return PPT_AGENT_STRATEGY_DEFAULT


def ppt_agent_strategy_label(strategy):
    strategy = normalize_ppt_agent_strategy(strategy)
    if strategy == PPT_AGENT_STRATEGY_DEFAULT:
        return "默认 PPT Agent"
    capability = PPT_HTML_CAPABILITIES.get(strategy)
    return capability.name if capability else "自动选择"


def ppt_agent_strategy_skill_name(strategy):
    strategy = normalize_ppt_agent_strategy(strategy)
    capability = PPT_HTML_CAPABILITIES.get(strategy)
    return capability.skill_name if capability else ""


def ppt_agent_builtin_skill_names():
    return [capability.skill_name for capability in PPT_HTML_CAPABILITIES.values()]


def validate_pptx_structure(path):
    """Validate PPTX structure without claiming to perform visual QA."""
    pptx_path = os.path.abspath(str(path or "").strip())
    if not pptx_path or not os.path.isfile(pptx_path):
        raise ValueError("PPTX 文件不存在。")
    if os.path.splitext(pptx_path)[1].lower() != ".pptx":
        raise ValueError("结构校验仅支持 PPTX 文件。")

    try:
        import pptx
    except Exception as exc:
        raise RuntimeError(f"python-pptx 不可用：{exc}") from exc

    required_members = {"[Content_Types].xml", "ppt/presentation.xml"}
    with zipfile.ZipFile(pptx_path, "r") as archive:
        members = set(archive.namelist())
        missing_members = sorted(required_members - members)
        if missing_members:
            raise ValueError("PPTX 缺少必要结构：" + "、".join(missing_members))
        corrupt_member = archive.testzip()
        if corrupt_member:
            raise ValueError(f"PPTX ZIP 成员损坏：{corrupt_member}")
        for rels_name in sorted(name for name in members if name.endswith(".rels")):
            try:
                rels_root = ElementTree.fromstring(archive.read(rels_name))
            except Exception as exc:
                raise ValueError(f"PPTX 关系文件无法解析：{rels_name}: {exc}") from exc
            rels_dir = posixpath.dirname(rels_name)
            source_dir = posixpath.dirname(rels_dir) if posixpath.basename(rels_dir) == "_rels" else rels_dir
            for relation in list(rels_root):
                if str(relation.attrib.get("TargetMode") or "").lower() == "external":
                    continue
                target = unquote(str(relation.attrib.get("Target") or "").split("#", 1)[0])
                if not target:
                    continue
                parsed = urlparse(target)
                if parsed.scheme:
                    continue
                normalized = (
                    target.lstrip("/")
                    if target.startswith("/")
                    else posixpath.normpath(posixpath.join(source_dir, target))
                )
                if normalized not in members:
                    raise ValueError(f"PPTX 关系目标不存在：{rels_name} -> {target}")

    try:
        presentation = pptx.Presentation(pptx_path)
    except Exception as exc:
        raise ValueError(f"python-pptx 无法重新打开成品：{exc}") from exc
    slide_width = int(presentation.slide_width or 0)
    slide_height = int(presentation.slide_height or 0)
    if slide_width <= 0 or slide_height <= 0:
        raise ValueError("PPTX 页面尺寸无效。")
    if len(presentation.slides) <= 0:
        raise ValueError("PPTX 没有任何幻灯片。")

    shape_count = 0
    text_shape_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            shape_count += 1
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            if width < 0 or height < 0:
                raise ValueError(f"第 {slide_index} 页存在负尺寸形状。")
            if bool(getattr(shape, "has_text_frame", False)):
                text_shape_count += 1
                _ = shape.text_frame.text
                if left < 0 or top < 0 or left + width > slide_width or top + height > slide_height:
                    raise ValueError(f"第 {slide_index} 页存在越出页面边界的文本框。")

    return {
        "path": pptx_path,
        "page_count": len(presentation.slides),
        "slide_width": slide_width,
        "slide_height": slide_height,
        "shape_count": shape_count,
        "text_shape_count": text_shape_count,
        "verification_level": "structure_only",
    }


def ppt_agent_capability_prompt_lines():
    lines = []
    for capability in PPT_HTML_CAPABILITIES.values():
        lines.extend(
            [
                f"- {capability.name}（已内置 Skill: {capability.skill_name}）: {capability.summary}",
                f"  适合: {capability.suitable_for}",
                f"  输出: {capability.output}，必须注册为 HTML deliverable，再复用 HTML→PPTX/DOCX/PDF 转换链路。",
            ]
        )
    return lines


def build_ppt_agent_prompt(
    request_text,
    preference=PPT_AGENT_PREFERENCE_AUTO,
    explicit_strategy=PPT_AGENT_STRATEGY_AUTO,
    template_file="",
    output_format=PPT_AGENT_OUTPUT_HTML,
    template_screenshots=None,
    renderer="",
    visual_validation=True,
):
    request = str(request_text or "").strip()
    output_format = normalize_ppt_agent_output_format(output_format)
    preference = normalize_ppt_agent_preference(preference)
    explicit_strategy = normalize_ppt_agent_strategy(explicit_strategy)
    template_file = str(template_file or "").strip()
    template_screenshots = [
        str(path or "").strip()
        for path in (template_screenshots or [])
        if str(path or "").strip()
    ]
    if output_format == PPT_AGENT_OUTPUT_PPTX:
        screenshot_lines = "\n".join(
            f"- 模板第 {index} 页截图: {path}"
            for index, path in enumerate(template_screenshots, start=1)
        )
        if not screenshot_lines:
            screenshot_lines = "- 当前没有本地演示渲染器，无法提供模板真实截图。"
        validation_line = (
            "- 生成后系统会用同一个本地演示渲染器（本机 PowerPoint 或 WPS）打开成品并逐页截图，交给多模态模型继续视觉校验；发现问题会自动发起修复轮次。"
            if visual_validation
            else "- 当前任务没有可用的 PowerPoint/WPS 渲染器；请完成结构校验，并明确说明成品未经视觉渲染校验。"
        )
        renderer_name = str(renderer or "none").strip()
        renderer_prog_ids = {
            "powerpoint": "PowerPoint.Application",
            "wps": "KWPP.Application / WPP.Application",
        }
        prog_id = renderer_prog_ids.get(renderer_name.lower())
        if prog_id:
            renderer_line = f"- 本地渲染器: {renderer_name}（ProgID: {prog_id}）——可自行复用\n"
        else:
            renderer_line = f"- 本地渲染器: {renderer_name or 'none'}\n"
        prompt = (
            "请以「PPT Agent」身份，根据用户材料和 PPTX 模板原文件直接生成新的 PPTX。"
            "本任务不要先生成 HTML，也不要使用默认 PPT Agent、Guizang PPT、Frontend Slides 或 Huashu Design 的 HTML-first 工作流。\n\n"
            "建议工作方式:\n"
            "1. 先查看本轮附加的全部模板截图，理解模板的视觉语言、页面类型、品牌元素和内容安全区。模板仅提供品牌规范、背景、色彩、页眉页脚和版式语言；模板中的示例文案、示例图形、参考流程、金字塔、图标和占位符不需要作为默认保留内容。\n"
            "2. 使用 python-pptx 程序化读取模板的页面尺寸、每个形状的坐标和大小、字体字号颜色、占位符、图片资源关系、母版和版式。\n"
            "3. python-pptx 无法覆盖的结构，可以直接检查 PPTX ZIP 包内的 OOXML、关系文件和媒体资源。\n"
            "4. 每页先区分固定保留元素、可复用版式容器和必须删除的示例内容元素。除固定品牌元素和明确复用的容器外，不得在原示例元素上直接叠加新内容。\n"
            "5. 先读取 Logo、公司名称等品牌元素坐标并设置标题安全区，所有标题必须避开品牌区域。模板页不适配内容时，使用保留品牌框架的空白页重建。\n"
            "6. 输出独立的新 PPTX，只保留生成页，不要覆盖或修改模板原文件。\n"
            "7. 完成后使用 python-pptx 重新打开生成文件，并检查 ZIP/OOXML、幻灯片数量、页面尺寸、形状坐标、文本边界、关系目标和媒体资源。\n"
            "8. 如本地有可用渲染器，必须通过 COM 自动化驱动本机 PowerPoint 或 WPS 实际打开成品并逐页导出全部截图。文件结构正确、能打开、页数正确均不能代替视觉验收。发现模板文字残留、旧图形与新内容重叠、标题侵入品牌区、文本截断或元素遮挡时，必须修复同一个 PPTX 后重新渲染并检查。\n\n"
            "模板与渲染信息:\n"
            f"- PPTX 模板原文件: {template_file}\n"
            f"{renderer_line}"
            f"{validation_line}\n"
            f"{screenshot_lines}\n\n"
            "[用户需求]\n"
            f"{request}"
        )
        return {
            "prompt": prompt,
            "selected_strategy": PPT_AGENT_STRATEGY_DEFAULT,
            "selected_label": "PPTX 模板驱动生成",
            "output_format": output_format,
        }
    selected_strategy = choose_ppt_agent_strategy(
        request,
        preference=preference,
        explicit_strategy=explicit_strategy,
        template_file=template_file,
    )
    selected_label = ppt_agent_strategy_label(selected_strategy)
    capability_lines = "\n".join(ppt_agent_capability_prompt_lines())
    template_lines = ""
    if template_file:
        template_lines = (
            "\n用户提供了 PPTX 模板，后续导出 PPTX 时必须把 HTML 作为内容来源、模板作为视觉与结构来源，"
            "不要修改原模板文件。\n"
            f"- PPTX 模板: {template_file}\n"
        )
    explicit_note = "用户显式指定了内置能力，优先尊重该选择。" if explicit_strategy != PPT_AGENT_STRATEGY_AUTO else "请根据需求自动选择最合适的生成策略。"
    prompt = (
        "请以「PPT Agent」身份处理下面的演示文稿任务，并生成可预览、可继续转换的 HTML 演示文稿工作稿。\n\n"
        f"生成偏好: {preference}\n"
        f"策略选择: {selected_label}\n"
        f"{explicit_note}\n\n"
        "PPT Agent 职责:\n"
        "- 判断用户要做的 PPT 类型，例如汇报、课件、路演、方案、研究报告、产品发布或视觉演示。\n"
        "- 先生成大纲和页面规划，再落地为演示文稿形态的完整 HTML 文件。\n"
        "- 不要只输出普通文本、Markdown 摘要或长文网页。\n"
        "- 不要绕过现有交付物系统直接生成 PPTX；PPTX/DOCX/PDF 后续统一从 HTML 工作稿转换。\n"
        "- 完成后明确给出项目工作区内 HTML 文件的完整路径，便于右侧交付物视图自动预览。\n\n"
        "HTML 工作稿约束:\n"
        "- 默认 16:9，按 slide 分页，每页有明确标题层级。\n"
        "- 内容适合投影展示，控制文字密度，保留图片、图表、布局和视觉层级信息。\n"
        "- 如果使用 CSS/JS/图片资源，请放在可迁移的本地目录，并确保 HTML deliverable preview 能正常打开。\n"
        "- HTML 应让后续转换 PPTX 时不需要重新理解整篇长文。\n\n"
        "已内置 html-ppt Skill:\n"
        f"{capability_lines}\n\n"
        "自动选择规则:\n"
        "- 普通汇报、课件、研究报告、会议总结: 默认 PPT Agent。\n"
        "- Word、PDF、Markdown 或长文本资料: 先整理内容结构，再生成 HTML 工作稿。\n"
        "- 提供 PPTX 模板: 优先生成符合模板转换要求的 HTML 工作稿，后续走模板化导出。\n"
        "- 网页 PPT、横向翻页、封面图、强风格内容分享: Guizang PPT Skill。\n"
        "- 产品发布、技术分享、交互式演示、前端风格 slides: Frontend Slides。\n"
        "- 高级感、商业汇报、发布会、路演、品牌提案、高审美: Huashu Design。\n"
        "- 如果选中的内置 Skill 不可用，请清晰告知并停止，不要静默降级。\n"
        f"{template_lines}\n"
        "[用户需求]\n"
        f"{request}"
    )
    return {
        "prompt": prompt,
        "selected_strategy": selected_strategy,
        "selected_label": selected_label,
    }
