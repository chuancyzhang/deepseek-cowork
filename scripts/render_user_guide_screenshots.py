import os
import sys
import tempfile
from dataclasses import replace
from pathlib import Path


os.environ.setdefault("QT_SCALE_FACTOR", "1")
os.environ.setdefault("QT_OPENGL", "software")
os.environ.setdefault("QTWEBENGINE_CHROMIUM_FLAGS", "--disable-gpu --disable-gpu-compositing")

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

TEMP_ROOT = Path(tempfile.mkdtemp(prefix="cowork-user-guide-"))
os.environ["APPDATA"] = str(TEMP_ROOT / "appdata")
APP_DATA_DIR = TEMP_ROOT / "appdata" / "DeepSeekCowork"
APP_DATA_DIR.mkdir(parents=True, exist_ok=True)
(APP_DATA_DIR / "config.json").write_text("{}", encoding="utf-8")

from PySide6.QtTest import QTest
from PySide6.QtCore import QPoint, Qt
from PySide6.QtGui import QImage
from PySide6.QtWidgets import QApplication, QMainWindow
from PySide6.QtWebEngineWidgets import QWebEngineView  # noqa: F401

import main


OUTPUT_DIR = Path(os.environ.get("COWORK_SCREENSHOT_OUTPUT_DIR") or (ROOT / "images" / "user-guide"))
SCREENSHOT_SCOPE = str(os.environ.get("COWORK_SCREENSHOT_SCOPE") or "").strip().lower()


def process_events(delay_ms=120):
    QApplication.processEvents()
    QTest.qWait(delay_ms)
    QApplication.processEvents()


def save_widget(widget, filename, delay_ms=120):
    widget.show()
    widget.raise_()
    process_events(delay_ms)
    pixmap = widget.grab()
    if pixmap.isNull():
        raise RuntimeError(f"Unable to capture screenshot: {filename}")
    output_path = OUTPUT_DIR / filename
    if not pixmap.save(str(output_path), "PNG"):
        raise RuntimeError(f"Unable to save screenshot: {output_path}")


def save_window_with_popups(window, filename, delay_ms=120):
    process_events(delay_ms)
    screen = window.screen() or QApplication.primaryScreen()
    geometry = window.frameGeometry()
    pixmap = screen.grabWindow(0, geometry.x(), geometry.y(), geometry.width(), geometry.height())
    output_path = OUTPUT_DIR / filename
    if pixmap.isNull() or not pixmap.save(str(output_path), "PNG"):
        raise RuntimeError(f"Unable to capture screenshot with popups: {output_path}")


def render_assistant_turn_screens(window):
    state = window.get_current_session()
    window._retire_session_empty_state(state, reason="user_guide_assistant_turn")
    window.add_chat_bubble(
        "User",
        "请检查工作区和现有测试，再给出最终修改说明。",
        animate=False,
        source_message_id="guide-turn-user",
    )
    turn_group = main.AssistantTurnGroup("guide-turn-group")
    state.chat_layout.insertWidget(state.chat_layout.count() - 1, turn_group)
    stage = window._create_agent_chat_bubble(state)
    window._connect_chat_bubble_actions(stage, state)
    turn_group.add_stage(stage, "guide-turn-group:stage-1")
    stage.update_thinking("先分析任务目标，再检查工作区文件。")
    tool = main.ToolCallCard("run_command", {"command": "python -m pytest"}, "guide-turn-tool")
    stage.add_tool_card(tool)
    tool.set_result("测试通过")
    stage.update_thinking(duration=10.2, is_final=True)
    stage.set_message_actions_enabled(False)
    stage.set_main_content("先完成环境与测试检查。", final=True)
    stage.think_toggle_btn.setChecked(True)
    final = window._create_agent_chat_bubble(state)
    window._connect_chat_bubble_actions(final, state)
    turn_group.add_stage(final, "guide-turn-group:stage-2")
    final.update_thinking("根据检查结果整理最终答复。")
    final.update_thinking(duration=4.1, is_final=True)
    final.set_source_message_id("guide-turn-final")
    final.set_main_content("现有测试已经通过。我已按验证结果整理最终修改说明。", final=True)
    save_widget(window, "37-thinking-expanded.png", 220)

    stage.think_toggle_btn.setChecked(False)
    final.set_message_actions_enabled(False)
    guidance = window.add_turn_guidance_inline(
        {"id": "guide-timeline-demo", "content": "先验证现有测试，再继续修改界面。"},
        status="waiting_tool",
    )
    save_widget(window, "s18-guidance-running.png", 220)
    if not isinstance(guidance, main.ChatBubble) or guidance.role != "User":
        raise RuntimeError("Running guidance was not rendered as a normal user message.")

    followup_group = main.AssistantTurnGroup("guide-followup-group")
    state.chat_layout.insertWidget(state.chat_layout.count() - 1, followup_group)
    followup = window._create_agent_chat_bubble(state)
    window._connect_chat_bubble_actions(followup, state)
    followup_group.add_stage(followup, "guide-followup-group:stage-1")
    followup.update_thinking("测试通过，继续整理最终结果。")
    followup.update_thinking(duration=4.1, is_final=True)
    followup.set_source_message_id("guide-followup-final")
    followup.set_main_content("现有测试已经通过。我已按验证结果整理最终修改说明。", final=True)
    save_widget(window, "s18-guidance-applied.png", 220)
    save_widget(window, "38-guidance-timeline.png", 220)


def render_history_performance_screens(window):
    state = window.get_current_session()
    window.clear_chat_layout(state.chat_layout)
    state.empty_state = None
    messages = [
        {
            "id": "history-guide-user",
            "role": "user",
            "content": "请分析这份长报告，并给出可以直接执行的优化方案。",
        }
    ]
    group_id = "history-guide-group"
    for index in range(6):
        messages.extend(
            [
                {
                    "id": f"history-stage-{index}",
                    "role": "assistant",
                    "content": f"阶段 {index + 1} 已完成。",
                    "reasoning_content": "正在核对数据、结构与约束。",
                    "tool_calls": [
                        {
                            "id": f"history-tool-{index}",
                            "function": {"name": "analyze_document", "arguments": "{}"},
                        }
                    ],
                    "meta": {
                        "ui_turn_group_id": group_id,
                        "ui_stage_id": f"{group_id}:{index}",
                        "ui_reply_kind": "stage",
                    },
                },
                {
                    "id": f"history-result-{index}",
                    "role": "tool",
                    "tool_call_id": f"history-tool-{index}",
                    "content": "检查完成",
                },
            ]
        )
    messages.append(
        {
            "id": "history-guide-final",
            "role": "assistant",
            "content": (
                "## 优化结论\n\n"
                "长历史已经改为**答案优先、执行详情按需加载**。"
                "打开会话时先呈现最终结果，需要时再展开完整过程。"
            ),
            "meta": {
                "ui_turn_group_id": group_id,
                "ui_stage_id": f"{group_id}:final",
                "ui_reply_kind": "final",
            },
        }
    )
    state.messages = messages
    state.render_items = main.build_conversation_render_spans(messages)
    state.displayed_render_count = len(state.render_items)
    window._render_session_history_spans(state, state.render_items)
    save_widget(window, "s37-history-on-demand.png", 260)

    window.component_task_manager._component_statuses = {
        "browser-skill": {
            "known": True,
            "installed": True,
            "state_text": "已安装，等待用户检查连接",
            "updated_at": 1784563200,
        },
        "node": {
            "known": True,
            "installed": True,
            "version": "v24.14.1",
            "updated_at": 1784563200,
        },
        "documents": {
            "known": True,
            "installed": True,
            "healthy": True,
            "size": 128 * 1024 * 1024,
            "updated_at": 1784563200,
        },
    }
    original_start_app_update = main.SettingsDialog.start_app_update
    main.SettingsDialog.start_app_update = lambda *_args, **_kwargs: None
    try:
        window.open_settings("组件与依赖")
    finally:
        main.SettingsDialog.start_app_update = original_start_app_update
    settings = window.product_pages[window.PAGE_SETTINGS]
    settings._automatic_update_check_started = True
    settings.components_scroll_area.ensureWidgetVisible(settings.browser_skill_component_section, 24, 24)
    save_widget(window, "s38-component-status-cache.png", 260)


def render_file_workbench_screens(window, workspace):
    window.load_workspace(
        str(workspace),
        refresh_sidebar=False,
        remember_workspace=False,
        persist_default=False,
        bind_session=True,
    )
    paths = [
        workspace / "weekly-brief.html",
        workspace / "meeting-notes.md",
        workspace / "release-notes.txt",
        workspace / "quarterly-review.pptx",
        workspace / "project-summary.pdf",
    ]
    for path in paths:
        window.chat_storage.register_deliverable(
            str(workspace),
            str(path),
            conversation_id=window.current_session_id,
            source="generated",
        )
    window.resize(1280, 760)
    window.show()
    window.context_drawer_expanded = False
    window.context_drawer_user_width = 470
    window.show_context_drawer(window.RIGHT_TAB_FILES)
    window.set_file_navigator_scope(
        window.FILE_SCOPE_DELIVERABLES,
        refresh=True,
        user_initiated=True,
    )
    window.file_navigator_pinned = False
    window.set_file_navigator_visible(True, reason="screenshot_list")
    for _attempt in range(160):
        process_events(20)
        if getattr(window, "deliverable_scan_worker", None) is None:
            break
    window.sync_context_drawer_layout()
    process_events(220)
    scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
    save_widget(window, f"file-workbench-narrow-{scale_label}x.png", 220)
    if scale_label == "1":
        save_widget(window, "s21-deliverables-list.png", 120)

    for tab_path in (
        workspace / "weekly-brief.html",
        workspace / "meeting-notes.md",
        workspace / "quarterly-review.pptx",
        workspace / "release-notes.txt",
    ):
        window.select_deliverable(str(tab_path), render_html=True)
    window.context_drawer_expanded = True
    window.context_drawer_user_width = 820
    window.file_navigator_pinned = True
    window.file_navigator_visible = True
    window.sync_context_drawer_layout()
    window._sync_file_navigator_layout()
    process_events(300)
    if len(window.file_tab_strip.paths) < 4:
        raise RuntimeError("File tab strip did not retain all opened files")
    if not window.file_workbench.is_effectively_pinned():
        raise RuntimeError("Pinned file navigator did not move into the left column")
    save_widget(window, f"file-workbench-expanded-{scale_label}x.png", 220)
    if scale_label == "1":
        save_widget(window, "s22-deliverable-preview.png", 120)

    window.select_deliverable(str(workspace / "quarterly-review.pptx"), render_html=True)
    process_events(240)
    save_widget(window, f"file-workbench-readonly-{scale_label}x.png", 180)

    window.hide_context_drawer(reason="screenshot_question_navigator")
    state = window.get_current_session()
    window.clear_chat_layout(state.chat_layout)
    state.empty_state = None
    state.render_nodes = {}
    state.render_node_by_message_id = {}
    messages = []
    for index, (question, answer) in enumerate(
        (
            ("请先检查工作区中有哪些资料", "已完成目录检查，并按文件类型整理了可用资料。"),
            ("把会议记录整理成一份项目简报", "简报已整理为目标、进展、风险与下周行动四部分。"),
            ("再生成一个适合汇报的 HTML 版本", "HTML 工作稿已生成，可在文件工作台中继续预览或编辑。"),
            ("最后核对交付物和测试结果", "交付物路径与关键测试均已核对完成。"),
        ),
        start=1,
    ):
        messages.extend(
            [
                {"id": f"guide-question-{index}", "role": "user", "content": question},
                {
                    "id": f"guide-answer-{index}",
                    "role": "assistant",
                    "content": answer,
                    "meta": {"ui_reply_kind": "final"},
                },
            ]
        )
    state.messages = messages
    state.render_items = main.build_conversation_render_spans(messages)
    state.displayed_render_count = len(state.render_items)
    window._render_session_history_spans(state, state.render_items)
    window._sync_question_navigator(state.session_id)
    process_events(220)
    entries = window._question_navigator_entries_for_state(state)
    anchor = QPoint(window.question_navigator_rail.width() + 10, window.question_navigator_rail.height() // 2)
    window._show_question_navigator_preview(entries[1], anchor)
    process_events(160)
    save_widget(window, f"question-navigator-{scale_label}x.png", 220)
    if scale_label == "1":
        save_widget(window, "s19a-question-navigator.png", 120)


def select_settings_page(dialog, label):
    for row in range(dialog.nav_list.count()):
        if dialog.nav_list.item(row).text() == label:
            dialog.nav_list.setCurrentRow(row)
            process_events()
            return
    raise RuntimeError(f"Settings page not found: {label}")


def render_enterprise_message_screens(window, qa=False, simulate_many=False):
    if simulate_many:
        base_spec = main.IM_PROVIDER_SPECS[0]
        simulated_specs = tuple(
            replace(
                base_spec,
                provider_id=f"mock_{index}",
                title=f"模拟渠道 {index + 1}",
                subtitle=f"用于验证扩展列表 {index + 1}",
                required_keys=(),
            )
            for index in range(12)
        )
        main.IM_PROVIDER_SPECS = simulated_specs
        main.IM_PROVIDER_ORDER = tuple(
            spec.provider_id for spec in simulated_specs
        )
    window.resize(1280, 720)
    window.open_settings("企业消息")
    settings = window.product_pages[window.PAGE_SETTINGS]
    settings._automatic_update_check_started = True
    window.show()
    process_events(180)
    settings._ensure_im_master_detail_layout()
    process_events(80)
    scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
    if simulate_many:
        save_widget(
            window,
            f"qa-enterprise-12-channels-{scale_label}x.png",
            180,
        )
        return
    if qa:
        settings._im_gateway_status_timer.stop()
        save_widget(window, f"qa-enterprise-messages-{scale_label}x.png", 180)
        window.resize(1024, 640)
        process_events(160)
        settings._ensure_im_master_detail_layout()
        save_widget(
            window,
            f"qa-enterprise-messages-narrow-{scale_label}x.png",
            120,
        )
        window.resize(1280, 720)
        process_events(140)
        settings._select_im_provider("wechat")
        settings.im_provider_detail_badges["wechat"].setText("连接中")
        settings.im_provider_detail_badges["wechat"].set_tone("primary")
        settings.im_provider_statuses["wechat"].set_text(
            "正在连接微信…",
            "info",
        )
        save_widget(window, f"qa-enterprise-connecting-{scale_label}x.png")
        settings.im_provider_detail_badges["wechat"].setText("使用中")
        settings.im_provider_detail_badges["wechat"].set_tone("success")
        settings.im_provider_statuses["wechat"].set_text(
            "微信正在使用中，收到的文字和链接会交给默认主助手处理。",
            "success",
        )
        save_widget(window, f"qa-enterprise-connected-{scale_label}x.png")
        settings.im_provider_detail_badges["wechat"].setText("连接失败")
        settings.im_provider_detail_badges["wechat"].set_tone("error")
        settings.im_provider_statuses["wechat"].set_text(
            "连接失败：微信登录已过期，请重新扫码接入。",
            "error",
        )
        save_widget(window, f"qa-enterprise-error-{scale_label}x.png")
        settings._select_im_provider("dingtalk")
        settings.im_provider_advanced_toggles["dingtalk"].setChecked(True)
        save_widget(window, f"qa-dingtalk-advanced-{scale_label}x.png")
    else:
        save_widget(window, "s40-enterprise-messages.png", 180)

    original_start_worker = main.ChannelQrDialog._start_worker
    main.ChannelQrDialog._start_worker = lambda _dialog: None
    try:
        qr_dialog = main.ChannelQrDialog("wechat", parent=settings)
    finally:
        main.ChannelQrDialog._start_worker = original_start_worker
    try:
        qr_dialog._on_qr_ready("https://weixin.qq.com/q/cowork-user-guide", 300)
        qr_dialog._on_status_changed("scanned")
        if qa:
            save_widget(qr_dialog, f"qa-wechat-scanned-{scale_label}x.png")
            qr_dialog._on_verify_code_required("请输入手机微信显示的配对码")
            save_widget(qr_dialog, f"qa-wechat-verify-{scale_label}x.png")
            qr_dialog._expires_at = 0
            qr_dialog._update_countdown()
            save_widget(qr_dialog, f"qa-wechat-expired-{scale_label}x.png")
            qr_dialog._on_failed("网络连接失败，请检查网络后重新生成二维码。")
            save_widget(qr_dialog, f"qa-wechat-failed-{scale_label}x.png")
        else:
            save_widget(qr_dialog, "s41-wechat-scan.png")
    finally:
        qr_dialog._timer.stop()
        qr_dialog.close()



def verify_drawer_layout(window, width, height):
    window.resize(width, height)
    window.sync_context_drawer_layout()
    process_events(80)
    drawer_rect = window.right_sidebar.geometry()
    if drawer_rect.right() > window.main_container.width() or drawer_rect.width() < main.DesignTokens.drawer_min_width:
        raise RuntimeError(f"Invalid drawer geometry at {width}x{height}: {drawer_rect}")
    input_origin = window.input_card.mapTo(window.main_container, QPoint(0, 0))
    input_right = input_origin.x() + window.input_card.width()
    if input_right > drawer_rect.x() - window.context_drawer_gap:
        raise RuntimeError(
            f"Composer overlaps drawer at {width}x{height}: input_right={input_right}, drawer_x={drawer_rect.x()}"
        )


def quiet_show_event(window, event):
    QMainWindow.showEvent(window, event)


def render_browser_skill_setup(window, app):
    dark = str(os.environ.get("COWORK_BROWSER_SKILL_DARK") or "").strip().lower() in {
        "1", "true", "yes", "on",
    }
    if dark:
        from unittest.mock import patch

        from core.theme import ThemeRuntimeManager, default_design_tokens
        from core.theme_service import ThemeRepository

        repository = ThemeRepository(str(APP_DATA_DIR))
        theme_manager = ThemeRuntimeManager(app, repository)
        app.theme_manager = theme_manager
        window.theme_manager = theme_manager
        theme_manager.themeChanged.connect(window._apply_runtime_theme)
        theme_manager.previewStateChanged.connect(window._on_theme_preview_state)
        repository.write_preview(
            name="浏览器设置深色验收",
            overrides={
                "tokens": {
                    "primary": "#8b93ff",
                    "primary_hover": "#9ca3ff",
                    "primary_pressed": "#737be8",
                    "primary_soft": "#252a43",
                    "primary_focus": "#6972cc",
                    "bg_app": "#101116",
                    "bg_main": "#15171d",
                    "bg_secondary": "#1b1e26",
                    "bg_disabled": "#22252d",
                    "bg_hover": "#242832",
                    "bg_pressed": "#2a2f3a",
                    "bg_sidebar": "#12141a",
                    "bg_sidebar_hover": "#1c2029",
                    "bg_sidebar_selected": "#23283a",
                    "bg_chat": "#101116",
                    "management_bg": "#101116",
                    "management_panel_bg": "#181b22",
                    "text_primary": "#eceef3",
                    "text_secondary": "#b7bcc8",
                    "text_tertiary": "#858c9b",
                    "text_disabled": "#666d7a",
                    "sidebar_text": "#eceef3",
                    "sidebar_text_muted": "#aeb4c2",
                    "sidebar_border": "#292d36",
                    "chat_text": "#eceef3",
                    "chat_text_muted": "#b7bcc8",
                    "chat_border": "#292d36",
                    "icon_primary": "#eceef3",
                    "icon_secondary": "#aeb4c2",
                    "border": "#343946",
                    "border_strong": "#474d5b",
                    "border_subtle": "#292d36",
                    "separator": "#292d36",
                    "warning_text": "#f0b86e",
                },
            },
            default_tokens=default_design_tokens(),
            session_id="browser-skill-screenshot",
        )
        with patch(
            "core.theme.QFontDatabase.families",
            return_value=["Microsoft YaHei UI", "Consolas"],
        ):
            if not theme_manager.apply_repository_state(reason="browser_skill_screenshot"):
                raise RuntimeError(theme_manager.last_error)

    window.skill_manager.load_skills()
    window.skill_manager_ready = True
    window.component_task_manager._component_statuses[main.BROWSER_SKILL_COMPONENT_ID] = {
        "known": True,
        "installed": True,
        "healthy": True,
        "ready": False,
        "state": "extension_disconnected",
        "state_text": "本地支持已准备，扩展未连接",
        "version": "0.1.8",
        "expected_version": "0.1.8",
        "bundle_ready": True,
        "bundle_error": "",
        "bundled_cli_available": True,
        "bundled_extension_available": True,
        "expected_extension_version": "0.1.4",
        "extension_prepared": True,
        "extension_prepared_version": "0.1.4",
        "extension_path": r"C:\Users\用户\AppData\Roaming\DeepSeekCowork\runtime_sandbox\v1\components\browser-skill-extension\extension",
        "available_browsers": [
            {"id": "chrome", "name": "Google Chrome", "path": r"C:\Program Files\Google\Chrome\Application\chrome.exe", "extensions_url": "chrome://extensions/"},
            {"id": "edge", "name": "Microsoft Edge", "path": r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe", "extensions_url": "edge://extensions/"},
        ],
        "protocol_incompatible": False,
        "updated_at": 1785686400,
    }
    original_probe = window.component_task_manager.probe_component
    window.component_task_manager.probe_component = lambda _component_id: True
    try:
        window.open_skills_center()
        skills_page = window.product_pages[window.PAGE_CAPABILITIES]
        browser_skill = next(
            (
                skill
                for skill in skills_page._all_skills
                if str(skill.get("name") or "") == "browser-automation"
            ),
            None,
        )
        if browser_skill is None:
            raise RuntimeError("Browser automation capability fixture was not found.")
        window.show_capability_detail(browser_skill)
        detail = window.product_pages["capability_detail"]
        edge_index = detail.browser_choice_combo.findData("edge")
        if edge_index < 0:
            raise RuntimeError("Edge choice was not rendered in BrowserSkill setup.")
        detail.browser_choice_combo.setCurrentIndex(edge_index)
        window.resize(1120, 720)
        window.show()
        process_events(220)
        scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
        tone = "dark" if dark else "light"
        save_widget(
            window,
            f"browser-skill-offline-{tone}-{scale_label}x.png",
            220,
        )
        if not dark and scale_label == "1":
            save_widget(window, "s29-capability-settings.png", 80)
    finally:
        window.component_task_manager.probe_component = original_probe


def main_run():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    workspace = TEMP_ROOT / "Linear Demo Project"
    workspace.mkdir(parents=True, exist_ok=True)
    html_path = workspace / "weekly-brief.html"
    html_path.write_text(
        """<!doctype html><html><head><meta charset='utf-8'><style>
        body{font-family:Segoe UI,Microsoft YaHei,sans-serif;margin:0;background:#f7f7f8;color:#202124}
        main{max-width:960px;margin:40px auto;background:white;padding:42px;border:1px solid #e6e6e9;border-radius:12px}
        h1{font-size:34px;margin:0 0 10px}p{color:#5f6269}.grid{display:grid;grid-template-columns:repeat(3,1fr);gap:14px;margin-top:30px}
        article{background:#f4f4f5;padding:18px;border-radius:8px}strong{color:#5e6ad2;font-size:24px;display:block}
        </style></head><body><main><h1>本周项目简报</h1><p>将任务进度、关键结果和下周行动集中在一页中。</p>
        <section class='grid'><article><strong>12</strong>完成任务</article><article><strong>4</strong>新增交付物</article><article><strong>92%</strong>按期完成</article></section>
        </main></body></html>""",
        encoding="utf-8",
    )
    (workspace / "meeting-notes.md").write_text(
        "# 项目会议记录\n\n- 完成本周报告\n- 整理交付物\n- 确认下周计划\n",
        encoding="utf-8",
    )
    (workspace / "release-notes.txt").write_text(
        "文件工作台验收\n\n"
        "1. 导航支持交付物与工作区文件切换。\n"
        "2. 搜索保持可见，类型和排序集中在筛选菜单。\n"
        "3. 预览与编辑通过工具栏中的铅笔和眼睛按钮切换。\n"
        "4. 未保存修改在切换文件、会话或关闭抽屉前需要确认。\n",
        encoding="utf-8",
    )
    from pptx import Presentation

    pptx_path = workspace / "quarterly-review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "季度项目复盘"
    slide.placeholders[1].text = "进展：完成文件工作台重构\n风险：窄窗口下导航转为浮层\n下一步：完成缩放验收"
    presentation.save(str(pptx_path))
    pdf_path = workspace / "project-summary.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n% demo")

    app = QApplication.instance() or QApplication([])
    main.initialize_desktop_theme(app)
    original_show_event = main.MainWindow.showEvent
    main.MainWindow.showEvent = quiet_show_event
    window = None
    try:
        window = main.MainWindow()
        window.resize(1280, 720)
        if SCREENSHOT_SCOPE == "favorites":
            project_dir = TEMP_ROOT / "weekly-report-project"
            project_dir.mkdir(parents=True, exist_ok=True)
            window.config_manager.set_skill_enabled("document-reader", True)
            window.skill_manager.load_skills()
            window.skill_manager_ready = True
            window.config_manager.upsert_project(str(project_dir), name="产品周报")
            window.config_manager.set_favorites(
                [
                    {
                        "id": "fav-weekly-report",
                        "name": "产品周报",
                        "description": "汇总本周进展、风险与下周计划。",
                        "prompt": "读取项目中的本周资料，生成结构清晰的产品周报。",
                        "skill_names": ["document-reader"],
                        "execution_mode": "workspace",
                        "workspace_dir": str(project_dir),
                        "schedule": {
                            "enabled": True,
                            "prompt_mode": "inherit",
                            "schedule_type": "weekly",
                            "time_of_day": "09:00",
                            "weekdays": [0],
                        },
                    },
                    {
                        "id": "fav-research",
                        "name": "快速研究",
                        "description": "在独立聊天中直接开始主题研究。",
                        "prompt": "研究我接下来输入的主题，并给出来源清晰的结论。",
                        "skill_names": ["web-search"],
                        "execution_mode": "chat",
                    },
                    {
                        "id": "fav-visual",
                        "name": "数据可视化模式",
                        "description": "加载可视化能力，等待输入具体数据和目标。",
                        "skill_names": ["visualize"],
                        "execution_mode": "chat",
                    },
                ]
            )
            window.open_favorites()
            window.resize(1280, 760)
            save_widget(window, "s40-favorites-library.png", 220)
            window.show_favorite_editor(favorite_id="fav-weekly-report")
            editor = window.product_pages["favorite_editor"]
            editor.run_options_toggle.setChecked(False)
            save_widget(window, "s41-favorite-editor.png", 220)
            editor.run_options_toggle.setChecked(True)
            process_events(160)
            editor_scroll = editor.findChild(main.QScrollArea)
            editor_scroll.verticalScrollBar().setValue(editor_scroll.verticalScrollBar().maximum())
            process_events(120)
            save_widget(window, "s42-favorite-schedule.png", 180)
            if os.environ.get("COWORK_SCREENSHOT_NARROW") == "1":
                window.resize(760, 720)
                save_widget(window, "favorites-editor-narrow.png", 180)
                window.resize(1280, 760)

            window.show_conversation_page()
            session_id = window.create_new_session(
                title="产品周报",
                make_current=True,
                workspace_dir=str(project_dir),
            )
            state = window.get_session(session_id)
            window._set_favorite_task_origin(
                state,
                window.config_manager.get_favorite("fav-weekly-report"),
                "scheduler",
                1786323600,
            )
            user_message_id = window._new_message_id()
            state.messages = [
                {
                    "id": user_message_id,
                    "role": "user",
                    "content": "读取项目中的本周资料，生成结构清晰的产品周报。",
                },
                {
                    "id": window._new_message_id(),
                    "role": "assistant",
                    "content": "产品周报已经整理完成，包含本周进展、主要风险和下周计划。",
                },
            ]
            window.add_chat_bubble(
                "User",
                state.messages[0]["content"],
                animate=False,
                source_message_id=user_message_id,
                session_id=session_id,
            )
            agent_bubble = window._create_agent_chat_bubble(
                state,
                text=state.messages[1]["content"],
                thinking="正在读取项目资料并整理周报结构。",
                duration=8.6,
            )
            window._connect_chat_bubble_actions(agent_bubble, state)
            agent_bubble.apply_dynamic_widths(window.dynamic_message_width, window.dynamic_user_bubble_width)
            window._attach_live_agent_stage(state, agent_bubble)
            agent_bubble.set_main_content(state.messages[1]["content"], final=True)
            agent_bubble.update_thinking(is_final=True)
            window.update_conversation_header()
            save_widget(window, "favorites-scheduled-chat.png", 220)

            from unittest.mock import patch

            from core.theme import ThemeRuntimeManager, default_design_tokens
            from core.theme_service import ThemeRepository

            repository = ThemeRepository(str(APP_DATA_DIR / "favorites-theme-preview"))
            theme_manager = ThemeRuntimeManager(app, repository)
            app.theme_manager = theme_manager
            window.theme_manager = theme_manager
            theme_manager.themeChanged.connect(window._apply_runtime_theme)
            theme_manager.previewStateChanged.connect(window._on_theme_preview_state)
            window.open_favorites()
            repository.write_preview(
                name="常用深色预览",
                overrides={
                    "tokens": {
                        "primary": "#8b93ff",
                        "primary_soft": "#252a43",
                        "bg_app": "#101116",
                        "bg_main": "#15171d",
                        "bg_panel": "#181b22",
                        "bg_sidebar": "#12141a",
                        "bg_sidebar_hover": "#1c2029",
                        "bg_sidebar_selected": "#23283a",
                        "management_bg": "#101116",
                        "management_panel_bg": "#181b22",
                        "text_primary": "#eceef3",
                        "text_secondary": "#b7bcc8",
                        "text_tertiary": "#858c9b",
                        "sidebar_text": "#eceef3",
                        "sidebar_text_muted": "#aeb4c2",
                        "sidebar_border": "#292d36",
                        "border": "#343946",
                        "border_subtle": "#292d36",
                        "separator": "#292d36",
                    }
                },
                default_tokens=default_design_tokens(),
                session_id="favorites-screenshot",
            )
            with patch(
                "core.theme.QFontDatabase.families",
                return_value=["Microsoft YaHei UI", "Consolas"],
            ):
                if not theme_manager.apply_repository_state(reason="favorites_screenshot"):
                    raise RuntimeError(theme_manager.last_error)
            save_widget(window, "favorites-theme-preview.png", 220)
            with patch(
                "core.theme.QFontDatabase.families",
                return_value=["Microsoft YaHei UI", "Consolas"],
            ):
                if not theme_manager.restore_saved_theme(reason="favorites_screenshot_restore"):
                    raise RuntimeError(theme_manager.last_error)
            return
        if SCREENSHOT_SCOPE == "browser-skill":
            render_browser_skill_setup(window, app)
            return
        if SCREENSHOT_SCOPE == "attachment-preview":
            image_path = ROOT / "images" / "user-guide" / "s09-interface-overview.png"
            if not image_path.is_file():
                raise RuntimeError(f"Attachment preview fixture is missing: {image_path}")
            window.show_conversation_page()
            window._add_prompt_files([str(image_path)])
            window.input_field.setPlainText("请分析这张界面截图，指出需要优化的区域。")
            window.resize(1280, 720)
            window.show()
            process_events(260)
            chips = window.prompt_files_section.findChildren(main.FileChip)
            if len(chips) != 1 or not chips[0]._is_image:
                raise RuntimeError("Image attachment chip was not rendered.")
            save_widget(window, "s16-pasted-image.png", 180)
            chips[0]._open_image_preview(str(image_path))
            process_events(180)
            dialog = chips[0]._preview_dialog
            if not main._qt_object_alive(dialog):
                raise RuntimeError("Image preview dialog did not open.")
            dialog.resize(980, 700)
            process_events(160)
            if (
                dialog.image_label.pixmap() is None
                or dialog.image_label.pixmap().isNull()
                or dialog.zoom_label.text() == ""
            ):
                raise RuntimeError("Image preview content or zoom state is missing.")
            save_widget(dialog, "s16-image-preview.png", 180)
            if os.environ.get("COWORK_SCREENSHOT_NARROW") == "1":
                dialog.resize(520, 420)
                process_events(160)
                save_widget(dialog, "attachment-preview-narrow.png", 120)
            dialog.close()
            return
        if SCREENSHOT_SCOPE == "deliverable-editors":
            window.load_workspace(
                str(workspace),
                refresh_sidebar=False,
                remember_workspace=False,
                persist_default=False,
                bind_session=True,
            )
            markdown_path = workspace / "meeting-notes.md"
            window.chat_storage.register_deliverable(
                str(workspace),
                str(markdown_path),
                conversation_id=window.current_session_id,
                source="generated",
            )
            window.show_context_drawer(window.RIGHT_TAB_FILES)
            window.set_file_navigator_scope(
                window.FILE_SCOPE_DELIVERABLES,
                refresh=True,
                user_initiated=True,
            )
            window.set_file_navigator_visible(False, reason="screenshot_editor")
            window.select_deliverable(str(markdown_path), render_html=True)
            window.begin_deliverable_edit()
            for _attempt in range(200):
                process_events(25)
                if window.deliverable_edit_state == "ready":
                    break
            if window.deliverable_edit_state != "ready":
                raise RuntimeError(
                    f"Deliverable editor did not become ready: {window.deliverable_edit_state}"
                )
            window.deliverable_text_editor.appendPlainText("- 补充发布前验证")
            window.resize(1440, 900)
            window.show()
            window.context_drawer_user_width = 760
            window.sync_context_drawer_layout()
            process_events(260)
            if (
                not window.deliverable_edit_action_bar.isHidden()
                or window.deliverable_edit_save_btn.isHidden()
                or window.deliverable_edit_save_btn.width() <= 0
                or window.deliverable_mode_btn.isHidden()
                or window.preview_stack.currentWidget()
                is not window.deliverable_text_editor_container
            ):
                raise RuntimeError("Deliverable edit controls are clipped or not active.")
            scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
            save_widget(
                window,
                f"s22a-deliverable-edit-{scale_label}x.png",
                260,
            )
            if scale_label == "1":
                save_widget(window, "s22a-deliverable-edit.png", 160)

            window._set_deliverable_dirty(False)
            window._release_deliverable_edit_session(show_preview=False)
            from docx import Document

            docx_path = workspace / "weekly-report.docx"
            document = Document()
            document.add_heading("项目周报", level=1)
            document.add_paragraph("这里是可以直接修改的正文。")
            document.sections[0].header.paragraphs[0].text = "项目资料"
            document.sections[0].footer.paragraphs[0].text = "内部使用"
            document.save(docx_path)
            window.chat_storage.register_deliverable(
                str(workspace),
                str(docx_path),
                conversation_id=window.current_session_id,
                source="generated",
            )
            window.select_deliverable(str(docx_path), render_html=True)
            window.begin_deliverable_edit()
            for _attempt in range(240):
                process_events(25)
                if window.deliverable_edit_state in {"ready", "blocked", "failed"}:
                    break
            if window.deliverable_edit_state != "ready":
                raise RuntimeError(
                    f"DOCX editor did not become ready: {window.deliverable_edit_state}"
                )
            if (
                window.deliverable_mode_btn.isHidden()
                or window.preview_stack.currentWidget()
                is not window.deliverable_editor_web_view
                or "页眉页脚将原样保留"
                not in window.deliverable_edit_status_label.text()
            ):
                raise RuntimeError("DOCX edit entry, content, or preservation state is missing.")
            save_widget(
                window,
                f"s22b-deliverable-docx-edit-{scale_label}x.png",
                260,
            )
            if scale_label == "1":
                save_widget(window, "s22b-deliverable-docx-edit.png", 160)
            return
        if SCREENSHOT_SCOPE == "sidebar-history-pagination":
            project_dir = TEMP_ROOT / "sidebar-pagination-project"
            project_dir.mkdir(parents=True, exist_ok=True)
            project_path = str(project_dir)
            window.config_manager.upsert_project(project_path, name="客户研究项目")
            for index in range(12):
                window.chat_storage.save_conversation(
                    f"sidebar-project-page-{index}",
                    [
                        {
                            "id": f"sidebar-project-message-{index}",
                            "role": "user",
                            "content": f"项目历史 {index}",
                        }
                    ],
                    title=f"项目分析记录 {index + 1}",
                    status="completed",
                    meta={
                        "workspace_dir": project_path,
                        "workspace_source": "project",
                    },
                )
            for index in range(8):
                window.chat_storage.save_conversation(
                    f"sidebar-chat-page-{index}",
                    [
                        {
                            "id": f"sidebar-chat-message-{index}",
                            "role": "user",
                            "content": f"独立聊天历史 {index}",
                        }
                    ],
                    title=f"独立聊天记录 {index + 1}",
                    status="completed",
                    meta={
                        "workspace_dir": str(TEMP_ROOT / "conversation-workspaces" / str(index)),
                        "workspace_source": "chat",
                    },
                )
            window.project_preview_paths.add(project_path)
            window.refresh_history_list()
            window.resize(900, 700)
            window.show()
            window.main_splitter.setSizes([main.DesignTokens.sidebar_min_width, 696])
            process_events(240)
            disclosure_copy = [
                button.text().strip()
                for button in window.history_disclosure_buttons.values()
            ]
            if len(window.history_rows) != 8 or disclosure_copy.count("展开显示") != 2:
                raise RuntimeError(
                    "Sidebar pagination fixture is invalid: "
                    f"rows={len(window.history_rows)} disclosures={disclosure_copy}"
                )
            if any(
                button.text().strip() == "显示更多历史"
                for button in window.history_container.findChildren(main.QPushButton)
                if button.isVisible()
            ):
                raise RuntimeError("Global history disclosure is still visible.")
            scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
            save_widget(
                window.sidebar,
                f"sidebar-history-pagination-{scale_label}x.png",
                220,
            )
            return
        if SCREENSHOT_SCOPE == "sidebar-activity":
            project_dir = TEMP_ROOT / "sidebar-activity-project"
            project_dir.mkdir(parents=True, exist_ok=True)
            project_path = str(project_dir)
            window.config_manager.upsert_project(project_path)
            for index in range(6):
                window.chat_storage.save_conversation(
                    f"sidebar-history-{index}",
                    [
                        {
                            "id": f"sidebar-history-message-{index}",
                            "role": "user",
                            "content": f"历史会话 {index}",
                        }
                    ],
                    title=f"历史会话 {index}",
                    status="completed",
                    meta={
                        "workspace_dir": project_path,
                        "workspace_source": "project",
                    },
                )
            session_id = window.create_new_session(
                make_current=True,
                workspace_dir=project_path,
            )
            state = window.get_session(session_id)
            state.messages = [
                {
                    "id": "sidebar-activity-user",
                    "role": "user",
                    "content": "这是一个非常长的会话标题，用来验证后台运行状态始终留在侧栏可视范围内",
                }
            ]
            state.session_status = "running"
            state.live_activity = True
            window.chat_storage.save_conversation(
                session_id,
                state.messages,
                title="这是一个非常长的会话标题，用来验证后台运行状态始终留在侧栏可视范围内",
                status="running",
                meta={
                    "workspace_dir": project_path,
                    "workspace_source": "project",
                },
            )
            window.project_preview_paths.add(project_path)
            window.refresh_history_list()
            window.resize(900, 650)
            window.show()
            window.main_splitter.setSizes([main.DesignTokens.sidebar_min_width, 696])
            process_events(220)
            scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
            status = window.history_activity_statuses[session_id]
            row = window.history_rows[session_id]
            status_right = status.mapTo(
                window.history_scroll.viewport(),
                status.rect().topRight(),
            ).x()
            if (
                status.isHidden()
                or status.geometry().right() > row.rect().right()
                or status_right >= window.history_scroll.viewport().width()
            ):
                raise RuntimeError(
                    "Sidebar activity status is clipped: "
                    f"row={row.geometry().getRect()} status={status.geometry().getRect()} "
                    f"hidden={status.isHidden()} viewport_right={status_right}/"
                    f"{window.history_scroll.viewport().width()}"
                )
            save_widget(
                window.sidebar,
                f"sidebar-activity-min-width-{scale_label}x.png",
                220,
            )
            row._set_actions_visible(True)
            process_events(120)
            status_right = status.mapTo(
                window.history_scroll.viewport(),
                status.rect().topRight(),
            ).x()
            if (
                status.isHidden()
                or status.geometry().right() > row.rect().right()
                or status_right >= window.history_scroll.viewport().width()
            ):
                raise RuntimeError(
                    "Sidebar activity status is clipped on hover: "
                    f"row={row.geometry().getRect()} status={status.geometry().getRect()} "
                    f"hidden={status.isHidden()} viewport_right={status_right}/"
                    f"{window.history_scroll.viewport().width()}"
                )
            save_widget(
                window.sidebar,
                f"sidebar-activity-hover-{scale_label}x.png",
                180,
            )
            original_status_color = main.DesignTokens.status_running
            try:
                main.DesignTokens.status_running = "#3157c8"
                window._apply_runtime_theme()
                window.history_activity_statuses[session_id].refresh_theme()
                process_events(120)
                save_widget(
                    window.sidebar,
                    f"sidebar-activity-theme-{scale_label}x.png",
                    180,
                )
            finally:
                main.DesignTokens.status_running = original_status_color
                window._apply_runtime_theme()
            return
        if SCREENSHOT_SCOPE == "skill-capture":
            skill_analysis = main.ConversationSkillEvidenceDialog(
                {
                    "confidence": "high",
                    "task_goal": {
                        "text": "把周报整理流程转成可重复执行的能力。",
                        "source_message_ids": ["demo-user"],
                    },
                    "outcome": {
                        "text": "已完成数据收集、结构整理和结果校验。",
                        "source_message_ids": ["demo-agent"],
                    },
                    "reusable_patterns": [
                        {
                            "text": "先统一收集进展，再按结果、风险和下周行动组织内容。",
                            "source_message_ids": ["demo-user", "demo-agent"],
                        }
                    ],
                    "missing_evidence": [],
                    "privacy_findings": [{"kind": "workspace_path"}],
                    "resource_candidates": [
                        {
                            "id": "weekly-report-reference",
                            "kind": "reference",
                            "description": "保存稳定的周报栏目定义",
                            "source_message_ids": ["demo-agent"],
                        }
                    ],
                },
                [],
                parent=window,
            )
            save_widget(skill_analysis, "s32-skill-capture.png", 200)
            skill_analysis.resize(620, 560)
            save_widget(skill_analysis, "s32-skill-capture-small.png", 200)
            skill_analysis.hide()
            state = window.get_current_session()
            window._retire_session_empty_state(state, reason="skill_capture_guide")
            window.add_chat_bubble(
                "User",
                "把这段周报整理流程沉淀为可复用的 Skill。",
                animate=False,
                source_message_id="skill-guide-user",
            )
            window.add_chat_bubble(
                "Assistant",
                "已完成流程梳理和结果校验。",
                animate=False,
                source_message_id="skill-guide-assistant",
            )
            state.pending_conversation_skill_result = {
                "capture_id": "skill-guide-capture",
                "phase": "analyzing",
            }
            window._update_skill_capture_status_card(
                state,
                "正在分析复用价值",
                detail="已转到后台，可继续对话",
                running=True,
                ensure_visible=True,
            )
            window.resize(980, 640)
            save_widget(window, "s33-skill-capture-background.png", 220)
            state.pending_conversation_skill_result["phase"] = "draft_ready"
            window._update_skill_capture_status_card(
                state,
                "Skill 草稿已生成",
                detail="待确认保存",
                pending=True,
            )
            save_widget(window, "s34-skill-capture-ready.png", 220)
            return
        if SCREENSHOT_SCOPE == "theme-acceptance":
            from unittest.mock import patch

            from core.theme import ThemeRuntimeManager, default_design_tokens
            from core.theme_service import ThemeRepository

            repository = ThemeRepository(str(APP_DATA_DIR))
            theme_manager = ThemeRuntimeManager(app, repository)
            app.theme_manager = theme_manager
            window.theme_manager = theme_manager
            theme_manager.themeChanged.connect(window._apply_runtime_theme)
            theme_manager.previewStateChanged.connect(window._on_theme_preview_state)
            repository.write_preview(
                name="星图工作台",
                overrides={
                    "font_scale": 1.05,
                    "density": "compact",
                    "tokens": {
                        "primary": "#8b93ff",
                        "bg_app": "#101116",
                        "bg_main": "#15171d",
                        "bg_sidebar": "#12141a",
                        "sidebar_text": "#eceef3",
                        "sidebar_text_muted": "#aeb4c2",
                        "sidebar_border": "#292d36",
                        "bg_sidebar_hover": "#1c2029",
                        "bg_sidebar_selected": "#23283a",
                        "bg_chat": "#15171d",
                        "text_primary": "#eceef3",
                        "text_secondary": "#b7bcc8",
                        "text_tertiary": "#858c9b",
                        "composer_bg": "#181b22",
                        "composer_text": "#eceef3",
                        "right_sidebar_bg": "#171920",
                        "right_sidebar_text": "#eceef3",
                        "right_sidebar_text_muted": "#aeb4c2",
                        "right_sidebar_header_bg": "#1c1f27",
                        "management_bg": "#101116",
                        "management_panel_bg": "#181b22",
                        "overlay_bg": "#181b22",
                        "overlay_text": "#eceef3",
                        "preview_shell_bg": "#181b22",
                        "preview_shell_text": "#eceef3",
                        "border": "#343946",
                        "border_subtle": "#292d36",
                        "separator": "#292d36",
                        "chat_border": "#292d36",
                        "composer_border": "#343946",
                        "right_sidebar_border": "#292d36",
                    },
                },
                workspace_scene={
                    "attachment": "fixed",
                    "layers": [
                        {"type": "solid", "color": "#171a21"},
                        {
                            "type": "stripes", "color": "#8b93ff", "spacing": 22,
                            "line_width": 1, "angle": 35, "size": 8, "opacity": 0.10,
                        },
                        {
                            "type": "grid", "color": "#8b93ff", "spacing": 28,
                            "line_width": 1, "major_every": 4,
                            "major_color": "rgba(139,147,255,0.16)", "opacity": 0.10,
                        },
                    ],
                },
                surfaces={
                    "shell.left_sidebar": {
                        "material": {"kind": "tint", "color": "#11131a", "opacity": 0.88}
                    },
                    "home.hero": {
                        "material": {"kind": "transparent"}
                    },
                    "conversation.composer": {
                        "material": {"kind": "tint", "color": "#171a21", "opacity": 0.94}
                    },
                    "home.reminder": {
                        "material": {"kind": "tint", "color": "#f8f9ff", "opacity": 0.96}
                    },
                },
                components={
                    "home.title": {"style": {"foreground": "#f3f4ff", "font_size": 23, "font_weight": 700}},
                    "home.card.ppt": {"layout": {"row": 0, "column": 0, "column_span": 2}, "style": {"border_color": "#555e95", "border_width": 1}},
                    "home.card.finance": {"layout": {"row": 1, "column": 0}},
                    "home.card.data": {"visible": False},
                    "home.card.browser": {"layout": {"row": 1, "column": 1}},
                },
                content={
                    "brand.title": "DeepSeek Cowork · 星图",
                    "home.title": "今晚，从一件重要的事开始",
                    "home.card.ppt.title": "星图演示",
                    "home.card.ppt.description": "进入固定行为的 PPT Mode",
                    "home.card.finance.title": "星图研究",
                    "home.card.browser.title": "自动操作网页",
                    "home.reminder.title": "主题改变呈现，不改变能力边界",
                    "composer.placeholder": "描述目标，界面会保持动作与安全边界不变",
                },
                default_tokens=default_design_tokens(),
                session_id="theme-acceptance",
            )
            with patch(
                "core.theme.QFontDatabase.families",
                return_value=["Microsoft YaHei UI", "Consolas"],
            ):
                if not theme_manager.apply_repository_state(reason="screenshot_acceptance"):
                    raise RuntimeError(theme_manager.last_error)
            window.resize(1280, 800)
            window.show()
            process_events(260)
            window.sync_context_drawer_layout()
            process_events(120)
            save_window_with_popups(window, "s39-ai-theme-package-home.png", 220)
            window.add_chat_bubble(
                "User",
                "请展示左右侧栏、聊天区、工具卡和输入区的完整主题覆盖。",
                animate=False,
                source_message_id="theme-acceptance-user",
            )
            state = window.get_current_session()
            group = main.AssistantTurnGroup("theme-acceptance-group")
            state.chat_layout.insertWidget(state.chat_layout.count() - 1, group)
            bubble = window._create_agent_chat_bubble(state)
            group.add_stage(bubble, "theme-acceptance-stage")
            bubble.update_thinking("正在检查主题绑定、Markdown CSS 与自绘控件。")
            tool = main.ToolCallCard(
                "validate_ui_theme",
                {"areas": ["sidebar", "conversation", "composer", "drawer"]},
                "theme-acceptance-tool",
            )
            bubble.add_tool_card(tool)
            tool.set_result("全部区域通过")
            bubble.update_thinking(duration=2.4, is_final=True)
            bubble.set_source_message_id("theme-acceptance-assistant")
            bubble.set_main_content(
                "## 全界面主题已应用\n\n- 左右侧栏与聊天画布\n- 输入区、工具卡和 Markdown\n- 固定安全恢复条",
                final=True,
            )
            window.show_context_drawer(window.RIGHT_TAB_OBSERVABILITY)
            window.show()
            process_events(180)
            scale_label = str(os.environ.get("QT_SCALE_FACTOR") or "1").replace(".", "_")
            for width, height in ((1280, 720), (1440, 900), (1920, 1080)):
                window.resize(width, height)
                process_events(120)
                window.sync_context_drawer_layout()
                process_events(80)
                verify_drawer_layout(window, width, height)
                save_widget(
                    window,
                    f"theme-acceptance-{width}x{height}-{scale_label}x.png",
                    220,
                )
            return
        if SCREENSHOT_SCOPE == "model-settings":
            window.open_settings("模型与服务")
            settings = window.product_pages[window.PAGE_SETTINGS]
            settings._automatic_update_check_started = True
            save_widget(window, "06-model-service-settings.png")
            model_empty = main.ModelEditDialog("openai", parent=settings)
            model_empty.model_name_input.setText("gpt-5.6")
            save_widget(model_empty, "07-add-model.png")
            model_empty.hide()
            model_filled = main.ModelEditDialog(
                "openai",
                {
                    "display_name": "GPT-5.6 Sol",
                    "model_name": "gpt-5.6",
                    "api_protocol": "responses",
                    "supports_vision": True,
                    "thinking_enabled": True,
                    "reasoning_efforts": ["none", "low", "medium", "high", "xhigh", "max"],
                    "reasoning_effort": "medium",
                },
                parent=settings,
            )
            save_widget(model_filled, "08-model-configuration.png")
            model_filled.hide()
            return
        if SCREENSHOT_SCOPE in {
            "enterprise-messages",
            "enterprise-messages-qa",
            "enterprise-messages-qa-12",
        }:
            render_enterprise_message_screens(
                window,
                qa=SCREENSHOT_SCOPE != "enterprise-messages",
                simulate_many=SCREENSHOT_SCOPE == "enterprise-messages-qa-12",
            )
            return
        if SCREENSHOT_SCOPE == "assistant-turn":
            if os.environ.get("COWORK_SCREENSHOT_NARROW") == "1":
                window.resize(900, 650)
            render_assistant_turn_screens(window)
            return
        if SCREENSHOT_SCOPE == "history-performance":
            render_history_performance_screens(window)
            return
        if SCREENSHOT_SCOPE == "file-workbench":
            render_file_workbench_screens(window, workspace)
            return
        if SCREENSHOT_SCOPE == "grill-mode":
            window.resize(1440, 900)
            window.input_card.setMinimumWidth(720)
            window.show()
            process_events(180)
            window.conversation_column.setFixedWidth(720)
            window.session_tabs.setFixedWidth(720)
            window.input_card.setFixedWidth(720)
            process_events(120)
            window.show_prompt_tool_menu()
            save_window_with_popups(window, "s15-composer-add-menu.png", 200)
            window.composer_action_popover.close()
            window.toggle_grill_mode()
            window.input_field.setPlainText("请先压力测试这份产品方案，再由我确认是否执行。")
            save_widget(window, "s15a-grill-mode-armed.png", 200)
            return
        if SCREENSHOT_SCOPE == "home":
            narrow = os.environ.get("COWORK_SCREENSHOT_NARROW") == "1"
            if narrow:
                window.resize(900, 720)
            else:
                window.resize(1440, 900)
            window.show()
            process_events(180)
            window.sync_context_drawer_layout()
            process_events(120)
            if narrow:
                state = window.get_current_session()
                scrollbar = state.chat_scroll.verticalScrollBar()
                scrollbar.setValue(scrollbar.maximum())
                process_events(120)
            save_widget(
                window,
                "s04-home-and-settings-narrow.png" if narrow else "s04-home-and-settings.png",
                250,
            )
            return
        save_widget(window, "04-home-screen.png", 250)
        model_style = window.model_select_btn.styleSheet()
        window.model_select_btn.setStyleSheet(main.apple_tool_button_style(True))
        save_widget(window, "09-switch-model.png")
        window.model_select_btn.setStyleSheet(model_style)
        window.show_model_selector_popover()
        if getattr(window, "model_selector_popover", None) is not None:
            save_window_with_popups(window, "33-model-popover.png")
            window.model_selector_popover.close()
            process_events(80)
        window.show_prompt_tool_menu()
        save_window_with_popups(window, "34-composer-add-popover.png", 200)
        window.composer_action_popover.close()
        project_add_style = window.sidebar_add_project_btn.styleSheet()
        window.sidebar_add_project_btn.setStyleSheet(main.apple_sidebar_icon_button_style(True))
        save_widget(window, "12-add-project.png")
        window.sidebar_add_project_btn.setStyleSheet(project_add_style)

        window.open_settings("模型与服务")
        settings = window.product_pages[window.PAGE_SETTINGS]
        settings._automatic_update_check_started = True
        save_widget(window, "05-open-settings.png")
        save_widget(window, "06-model-service-settings.png")

        model_empty = main.ModelEditDialog("openai", parent=settings)
        save_widget(model_empty, "07-add-model.png")
        model_empty.hide()
        model_filled = main.ModelEditDialog(
            "openai",
            {
                "display_name": "DeepSeek V4 Pro",
                "model_name": "deepseek-v4-pro",
                "api_protocol": "chat_completions",
                "supports_vision": True,
                "thinking_enabled": True,
                "reasoning_efforts": ["low", "medium", "high"],
            },
            parent=settings,
        )
        save_widget(model_filled, "08-model-configuration.png")
        model_filled.hide()
        select_settings_page(settings, "更新")
        save_widget(window, "10-update-settings.png")
        select_settings_page(settings, "个性与记忆")
        save_widget(window, "25-memory-center.png")
        window.show_conversation_page()

        window.skill_manager.load_skills()
        window.skill_manager_ready = True
        window.open_session_skill_picker()
        if getattr(window, "session_skill_popover", None) is not None:
            save_window_with_popups(window, "35-ability-popover.png")
            window.session_skill_popover.close()
        window.open_skills_center()
        skills = window.product_pages[window.PAGE_CAPABILITIES]
        save_widget(window, "23-skills-center.png")
        browser_skill = next(
            (item for item in skills._all_skills if item.get("name") == "browser-automation"),
            None,
        )
        if browser_skill:
            window.show_capability_detail(browser_skill)
            save_widget(window, "s29-capability-settings.png")
        window.show_conversation_page()

        window.config_manager.upsert_favorite(
            {
                "id": "guide-favorite",
                "name": "产品周报",
                "prompt": "汇总本周产品进展并生成周报。",
                "execution_mode": "chat",
            }
        )
        window.open_favorites()
        save_widget(window, "s40-favorites-library.png")
        window.show_favorite_editor(favorite_id="guide-favorite")
        editor = window.product_pages["favorite_editor"]
        editor.run_options_toggle.setChecked(False)
        save_widget(window, "s41-favorite-editor.png")
        window.handle_product_back()
        window.show_conversation_page()

        agent_style = window.agent_picker_btn.styleSheet()
        window.agent_picker_btn.setStyleSheet(main.apple_button_style("selected", radius=7))
        save_widget(window, "26-agent-center.png")
        window.agent_picker_btn.setStyleSheet(agent_style)

        skill_analysis = main.ConversationSkillEvidenceDialog(
            {
                "confidence": "high",
                "task_goal": {
                    "text": "把周报整理流程转成可重复执行的能力。",
                    "source_message_ids": ["demo-user"],
                },
                "outcome": {
                    "text": "已完成数据收集、结构整理和结果校验。",
                    "source_message_ids": ["demo-agent"],
                },
                "reusable_patterns": [
                    {
                        "text": "先统一收集进展，再按结果、风险和下周行动组织内容。",
                        "source_message_ids": ["demo-user", "demo-agent"],
                    }
                ],
                "missing_evidence": [],
                "privacy_findings": [{"kind": "workspace_path"}],
                "resource_candidates": [
                    {
                        "id": "weekly-report-reference",
                        "kind": "reference",
                        "description": "保存稳定的周报栏目定义",
                        "source_message_ids": ["demo-agent"],
                    }
                ],
            },
            [],
            parent=window,
        )
        save_widget(skill_analysis, "31-skill-capture-analysis.png")
        skill_analysis.hide()

        ppt_agent = main.PptAgentModeDialog(str(workspace), parent=window)
        save_widget(ppt_agent, "27-ppt-agent.png")
        ppt_agent.hide()

        user_bubble = window.add_chat_bubble(
            "User",
            "请整理本周项目进展，并生成一份适合汇报的简报。",
            animate=False,
            source_message_id="guide-user-message",
        )
        window.add_chat_bubble(
            "Agent",
            "已整理任务进度、关键结果和下周行动。你可以继续让我生成 HTML 工作稿，或直接调整内容结构。",
            animate=False,
        )
        process_events(200)
        save_widget(window, "11-direct-chat.png")
        if user_bubble is not None:
            user_bubble.begin_inline_edit()
            save_widget(user_bubble, "36-history-message-edit.png", 200)
            user_bubble.cancel_inline_edit()
        state = window.get_current_session()
        turn_group = main.AssistantTurnGroup("guide-turn-group")
        state.chat_layout.insertWidget(state.chat_layout.count() - 1, turn_group)
        thinking_bubble = window._create_agent_chat_bubble(state)
        window._connect_chat_bubble_actions(thinking_bubble, state)
        turn_group.add_stage(thinking_bubble, "guide-turn-group:stage-1")
        if thinking_bubble is not None:
            thinking_bubble.update_thinking("先分析任务目标，再检查工作区文件。")
            demo_tool = main.ToolCallCard("run_command", {"command": "python -m unittest"}, "guide-demo-tool")
            thinking_bubble.add_tool_card(demo_tool)
            demo_tool.set_result("测试通过")
            thinking_bubble.update_thinking(duration=10.2, is_final=True)
            thinking_bubble.set_message_actions_enabled(False)
            thinking_bubble.set_main_content("先完成环境与测试检查。", final=True)
            thinking_bubble.think_toggle_btn.setChecked(True)

            final_bubble = window._create_agent_chat_bubble(state)
            window._connect_chat_bubble_actions(final_bubble, state)
            turn_group.add_stage(final_bubble, "guide-turn-group:stage-2")
            final_bubble.update_thinking("根据检查结果整理最终答复。")
            final_bubble.update_thinking(duration=4.1, is_final=True)
            final_bubble.set_source_message_id("guide-final-message")
            final_bubble.set_main_content("现有测试已经通过。我已按验证结果整理最终修改说明。", final=True)
            save_widget(window, "37-thinking-expanded.png")
            thinking_bubble.think_toggle_btn.setChecked(False)
            final_bubble.set_message_actions_enabled(False)

            window.add_turn_guidance_inline(
                {"id": "guide-timeline-demo", "content": "先验证现有测试，再继续修改界面。"},
                status="applied",
            )
            followup_group = main.AssistantTurnGroup("guide-followup-group")
            state.chat_layout.insertWidget(state.chat_layout.count() - 1, followup_group)
            followup_bubble = window._create_agent_chat_bubble(state)
            window._connect_chat_bubble_actions(followup_bubble, state)
            followup_group.add_stage(followup_bubble, "guide-followup-group:stage-1")
            followup_bubble.update_thinking("测试通过，继续整理最终结果。")
            followup_bubble.update_thinking(duration=4.1, is_final=True)
            followup_bubble.set_source_message_id("guide-followup-final")
            followup_bubble.set_main_content(
                "现有测试已经通过。我已按验证结果整理最终修改说明。",
                final=True,
            )
            save_widget(window, "38-guidance-timeline.png")

        clipboard_image = QImage(240, 140, QImage.Format_ARGB32)
        clipboard_image.fill(Qt.GlobalColor.lightGray)
        window._add_clipboard_image(clipboard_image)
        process_events(120)
        save_widget(window, "13-paste-image.png")
        window._clear_prompt_files(window.current_session_id)

        window.load_workspace(
            str(workspace),
            refresh_sidebar=False,
            remember_workspace=False,
            persist_default=False,
            bind_session=True,
        )
        process_events(200)
        save_widget(window, "14-project-workspace.png")
        window.add_chat_bubble(
            "Agent",
            f"HTML 工作稿已生成：{html_path.name}\n\n已按 16:9 汇报节奏整理，可以在右侧预览并继续转换。",
            animate=False,
        )
        window.chat_storage.register_deliverable(
            str(workspace),
            str(html_path),
            conversation_id=window.current_session_id,
            source="generated",
        )
        window.chat_storage.register_deliverable(
            str(workspace), str(pptx_path), conversation_id=window.current_session_id, source="converted"
        )
        window.chat_storage.register_deliverable(
            str(workspace), str(pdf_path), conversation_id=window.current_session_id, source="converted"
        )
        process_events(200)
        save_widget(window, "15-generate-html-example.png")

        window.show_context_drawer(window.RIGHT_TAB_FILES)
        window.set_file_navigator_scope(window.FILE_SCOPE_WORKSPACE, refresh=True, user_initiated=True)
        window.set_file_navigator_visible(True, reason="screenshot_workspace")
        process_events(500)
        save_widget(window, "16-open-deliverables.png")
        window.set_file_navigator_scope(window.FILE_SCOPE_DELIVERABLES, refresh=True, user_initiated=True)
        process_events(250)
        save_widget(window, "17-deliverables-panel.png")
        window.add_system_toast("PPTX 已生成，可在交付物中打开。", "success", auto_close_ms=0)
        window.add_system_toast("资源管理器已打开并定位到文件。", "info", auto_close_ms=0)
        process_events(120)
        save_widget(window, "30-system-feedback.png")
        for toast in list(window._visible_system_toasts):
            window._dismiss_system_toast(toast)

        window.set_file_navigator_visible(False, reason="screenshot_preview")
        window.select_deliverable(str(html_path), render_html=True)
        process_events(1200)
        save_widget(window, "18-open-html-preview.png")
        save_widget(window, "20-generate-pptx.png")
        tool_id = "guide-python-tool"
        window.add_tool_card(
            {
                "id": tool_id,
                "name": "run_python_code",
                "args": {
                    "code": "from pathlib import Path\n\nfiles = list(Path('.').glob('*.md'))\nprint({'files': len(files)})",
                    "timeout": 30,
                },
                "meta": {"start_time": 1783783680, "duration": 0.24},
            },
            session_id=window.current_session_id,
            animate=False,
        )
        window.update_tool_card(
            {
                "id": tool_id,
                "result": '{"files": 1, "status": "ok"}',
                "result_obj": {"files": 1, "status": "ok"},
                "meta": {"duration": 0.24},
            },
            session_id=window.current_session_id,
        )
        window.show_context_drawer(window.RIGHT_TAB_OBSERVABILITY)
        card = window.get_current_session().tool_cards[tool_id]
        window.show_tool_details(tool_id, card.args, card.result, meta=card.meta, switch_tab=True)
        process_events(180)
        save_widget(window, "28-task-observability.png")
        for width, height in ((1280, 720), (1440, 900), (1920, 1080)):
            verify_drawer_layout(window, width, height)
            save_widget(window, f"38-responsive-{width}x{height}.png", 80)
    finally:
        main.MainWindow.showEvent = original_show_event
        if window is not None:
            worker = getattr(window, "chat_save_worker", None)
            if worker is not None and worker.isRunning():
                worker.stop_worker(timeout_ms=3000)
            window.hide()
        app.quit()


if __name__ == "__main__":
    main_run()
