import json
import os

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from core.env_utils import ensure_package_installed
from core.filesystem_ops import (
    _build_error,
    _build_ok,
    delete_path,
    ensure_existing_file_write_allowed,
    glob_paths,
    grep_contents,
    list_files as list_files_core,
    mark_file_written,
    read_text_file,
    record_full_read_state,
    rename_path,
    resolve_path,
    update_text_file,
    write_text_file,
)
from core.interaction import ask_user


def _get_openpyxl():
    try:
        import openpyxl

        return openpyxl
    except ImportError:
        ensure_package_installed("openpyxl", skill_id="file-system")
        import openpyxl

        return openpyxl


def _build_read_payload(action, rel_path, content):
    lines = content.splitlines()
    return _build_ok(
        action,
        {
            "path": rel_path,
            "content": content,
            "encoding": "utf-8",
            "truncated": False,
            "start_line": 1,
            "returned_lines": len(lines),
            "total_lines": len(lines),
        },
    )


def _parse_pdf_pages(pages, total_pages, action, rel_path):
    if pages is None:
        return list(range(total_pages)), None

    text = str(pages).strip()
    if not text:
        return list(range(total_pages)), None

    if "-" in text:
        start_raw, end_raw = text.split("-", 1)
        try:
            start_page = int(start_raw)
            end_page = int(end_raw)
        except Exception:
            return None, _build_error(action, "invalid_pages", "pages must be in 'N' or 'N-M' format.", path=rel_path)
        if start_page <= 0 or end_page <= 0 or start_page > end_page:
            return None, _build_error(action, "invalid_pages", "Invalid page range.", path=rel_path)
        first = start_page - 1
        last = min(end_page, total_pages)
        if first >= total_pages:
            return None, _build_error(action, "invalid_pages", "Page range exceeds PDF page count.", path=rel_path)
        return list(range(first, last)), None

    try:
        single = int(text)
    except Exception:
        return None, _build_error(action, "invalid_pages", "pages must be in 'N' or 'N-M' format.", path=rel_path)
    if single <= 0 or single > total_pages:
        return None, _build_error(action, "invalid_pages", "Requested page is out of range.", path=rel_path)
    return [single - 1], None


def list_files(workspace_dir, path=".", recursive=False, include_hidden=False, limit=200, _context=None):
    return list_files_core(
        workspace_dir,
        path=path,
        recursive=recursive,
        include_hidden=include_hidden,
        limit=limit,
        context=_context,
    )


def read_file(workspace_dir, path, offset=1, limit=None, sheet_name=None, pages=None, _context=None):
    ext = os.path.splitext(str(path or ""))[1].lower()
    if ext == ".docx":
        return read_docx(workspace_dir, path, _context=_context)
    if ext == ".pptx":
        return read_pptx(workspace_dir, path, _context=_context)
    if ext == ".xlsx":
        return read_excel(workspace_dir, path, sheet_name=sheet_name, _context=_context)
    if ext == ".pdf":
        return read_pdf(workspace_dir, path, pages=pages, _context=_context)
    return read_text_file(workspace_dir, path, offset=offset, limit=limit, context=_context, action="read_file")


def write_file(workspace_dir, path, content, mode="overwrite", _context=None):
    return write_text_file(
        workspace_dir,
        path,
        content,
        mode=mode,
        context=_context,
        action="write_file",
    )


def update_file(workspace_dir, path, old_string, new_string, replace_all=False, _context=None):
    return update_text_file(
        workspace_dir,
        path,
        old_string=old_string,
        new_string=new_string,
        replace_all=replace_all,
        context=_context,
        action="update_file",
    )


def rename_file(workspace_dir, old_path, new_path, _context=None):
    return rename_path(
        workspace_dir,
        old_path,
        new_path,
        context=_context,
        action="rename_file",
    )


def delete_file(workspace_dir, path, recursive=False, _context=None):
    def _confirm(rel_path, recursive_flag):
        if recursive_flag:
            prompt = f"Confirm delete recursively: '{rel_path}'?"
        else:
            prompt = f"Confirm delete: '{rel_path}'?"
        return ask_user(prompt)

    return delete_path(
        workspace_dir,
        path,
        recursive=recursive,
        confirm_callback=_confirm,
        context=_context,
        action="delete_file",
    )


def glob(workspace_dir, pattern="*", path=".", limit=200, include_hidden=False, _context=None):
    return glob_paths(
        workspace_dir,
        pattern=pattern,
        path=path,
        limit=limit,
        include_hidden=include_hidden,
        context=_context,
    )


def grep(workspace_dir, pattern, path=".", include="*", exclude=None, recursive=True, limit=200, _context=None):
    return grep_contents(
        workspace_dir,
        pattern=pattern,
        path=path,
        include=include,
        exclude=exclude,
        recursive=recursive,
        limit=limit,
        context=_context,
    )


def read_docx(workspace_dir, path, _context=None):
    action = "read_docx"
    abs_path, rel_path, error = resolve_path(workspace_dir, path, context=_context, action=action, must_exist=True)
    if error:
        return error
    if not os.path.isfile(abs_path):
        return _build_error(action, "not_a_file", "Path is not a file.", path=rel_path)

    try:
        doc = Document(abs_path)
        content = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        record_full_read_state(abs_path, _context)
        return _build_read_payload(action, rel_path, content)
    except Exception as exc:
        return _build_error(action, "read_failed", str(exc), path=rel_path)


def write_docx(workspace_dir, path, content, mode="w", _context=None):
    action = "write_docx"
    abs_path, rel_path, error = resolve_path(
        workspace_dir,
        path,
        context=_context,
        action=action,
        must_exist=False,
        reject_glob_for_write=True,
    )
    if error:
        return error

    mode_value = str(mode or "w").strip().lower() or "w"
    if mode_value not in {"w", "a"}:
        return _build_error(action, "invalid_mode", "mode must be 'w' or 'a'.", path=rel_path)

    existed_before = os.path.exists(abs_path)
    if existed_before:
        error = ensure_existing_file_write_allowed(abs_path, rel_path, _context, action)
        if error:
            return error

    parent = os.path.dirname(abs_path)
    if parent:
        os.makedirs(parent, exist_ok=True)

    text = content if isinstance(content, str) else str(content)
    try:
        if mode_value == "a" and existed_before:
            doc = Document(abs_path)
        else:
            doc = Document()
        for line in text.split("\n"):
            doc.add_paragraph(line)
        doc.save(abs_path)
        mark_file_written(abs_path, _context)
        change_type = "create" if not existed_before else ("append" if mode_value == "a" else "update")
        return _build_ok(
            action,
            {"path": rel_path, "change_type": change_type, "bytes_written": len(text.encode("utf-8"))},
        )
    except Exception as exc:
        return _build_error(action, "write_failed", str(exc), path=rel_path)


def read_pptx(workspace_dir, path, _context=None):
    action = "read_pptx"
    abs_path, rel_path, error = resolve_path(workspace_dir, path, context=_context, action=action, must_exist=True)
    if error:
        return error
    if not os.path.isfile(abs_path):
        return _build_error(action, "not_a_file", "Path is not a file.", path=rel_path)

    try:
        presentation = Presentation(abs_path)
        chunks = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            chunks.append(f"Slide {index}:\n" + "\n".join(texts))
        content = "\n\n".join(chunks)
        record_full_read_state(abs_path, _context)
        return _build_read_payload(action, rel_path, content)
    except Exception as exc:
        return _build_error(action, "read_failed", str(exc), path=rel_path)


def create_pptx(workspace_dir, path, slides_data, _context=None):
    action = "create_pptx"
    abs_path, rel_path, error = resolve_path(
        workspace_dir,
        path,
        context=_context,
        action=action,
        must_exist=False,
        reject_glob_for_write=True,
    )
    if error:
        return error

    existed_before = os.path.exists(abs_path)
    if existed_before:
        error = ensure_existing_file_write_allowed(abs_path, rel_path, _context, action)
        if error:
            return error

    parent = os.path.dirname(abs_path)
    if parent:
        os.makedirs(parent, exist_ok=True)

    if isinstance(slides_data, str):
        try:
            slides_data = json.loads(slides_data)
        except Exception:
            return _build_error(action, "invalid_argument", "slides_data must be a JSON list or list object.", path=rel_path)
    if not isinstance(slides_data, list):
        return _build_error(action, "invalid_argument", "slides_data must be a list.", path=rel_path)

    try:
        presentation = Presentation()
        for slide_info in slides_data:
            if not isinstance(slide_info, dict):
                continue
            title_text = str(slide_info.get("title") or "")
            content_text = str(slide_info.get("content") or "")
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = content_text
        presentation.save(abs_path)
        serialized = json.dumps(slides_data, ensure_ascii=False)
        mark_file_written(abs_path, _context)
        change_type = "create" if not existed_before else "update"
        return _build_ok(
            action,
            {"path": rel_path, "change_type": change_type, "bytes_written": len(serialized.encode("utf-8"))},
        )
    except Exception as exc:
        return _build_error(action, "write_failed", str(exc), path=rel_path)


def read_excel(workspace_dir, path, sheet_name=None, _context=None):
    action = "read_excel"
    abs_path, rel_path, error = resolve_path(workspace_dir, path, context=_context, action=action, must_exist=True)
    if error:
        return error
    if not os.path.isfile(abs_path):
        return _build_error(action, "not_a_file", "Path is not a file.", path=rel_path)

    try:
        openpyxl = _get_openpyxl()
        workbook = openpyxl.load_workbook(abs_path, data_only=True)
        if sheet_name:
            if sheet_name not in workbook.sheetnames:
                return _build_error(action, "sheet_not_found", f"Sheet '{sheet_name}' not found.", path=rel_path)
            worksheet = workbook[sheet_name]
        else:
            worksheet = workbook.active

        rows = []
        for row in worksheet.iter_rows(values_only=True):
            rows.append("\t".join("" if cell is None else str(cell) for cell in row))
        content = "\n".join(rows)
        record_full_read_state(abs_path, _context)
        return _build_read_payload(action, rel_path, content)
    except Exception as exc:
        return _build_error(action, "read_failed", str(exc), path=rel_path)


def write_excel(workspace_dir, path, data, sheet_name="Sheet1", _context=None):
    action = "write_excel"
    abs_path, rel_path, error = resolve_path(
        workspace_dir,
        path,
        context=_context,
        action=action,
        must_exist=False,
        reject_glob_for_write=True,
    )
    if error:
        return error

    if isinstance(data, str):
        try:
            data = json.loads(data)
        except Exception:
            return _build_error(action, "invalid_argument", "data must be a JSON list of rows.", path=rel_path)
    if not isinstance(data, list):
        return _build_error(action, "invalid_argument", "data must be a list of rows.", path=rel_path)

    existed_before = os.path.exists(abs_path)
    if existed_before:
        error = ensure_existing_file_write_allowed(abs_path, rel_path, _context, action)
        if error:
            return error

    parent = os.path.dirname(abs_path)
    if parent:
        os.makedirs(parent, exist_ok=True)

    try:
        openpyxl = _get_openpyxl()
        if existed_before:
            workbook = openpyxl.load_workbook(abs_path)
            if sheet_name in workbook.sheetnames:
                del workbook[sheet_name]
            worksheet = workbook.create_sheet(sheet_name)
        else:
            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet.title = sheet_name

        for row in data:
            worksheet.append(row if isinstance(row, list) else [row])
        workbook.save(abs_path)
        serialized = json.dumps(data, ensure_ascii=False)
        mark_file_written(abs_path, _context)
        change_type = "create" if not existed_before else "update"
        return _build_ok(
            action,
            {"path": rel_path, "change_type": change_type, "bytes_written": len(serialized.encode("utf-8"))},
        )
    except Exception as exc:
        return _build_error(action, "write_failed", str(exc), path=rel_path)


def read_pdf(workspace_dir, path, pages=None, _context=None):
    action = "read_pdf"
    abs_path, rel_path, error = resolve_path(workspace_dir, path, context=_context, action=action, must_exist=True)
    if error:
        return error
    if not os.path.isfile(abs_path):
        return _build_error(action, "not_a_file", "Path is not a file.", path=rel_path)

    try:
        reader = PdfReader(abs_path)
        page_indices, error = _parse_pdf_pages(pages, len(reader.pages), action, rel_path)
        if error:
            return error
        blocks = []
        for page_index in page_indices:
            text = reader.pages[page_index].extract_text() or ""
            blocks.append(f"--- Page {page_index + 1} ---\n{text}")
        content = "\n\n".join(blocks)
        record_full_read_state(abs_path, _context)
        payload = {
            "path": rel_path,
            "content": content,
            "encoding": "utf-8",
            "truncated": False,
            "start_line": 1,
            "returned_lines": len(content.splitlines()),
            "total_lines": len(content.splitlines()),
            "page_count": len(page_indices),
        }
        return _build_ok(action, payload)
    except Exception as exc:
        return _build_error(action, "read_failed", str(exc), path=rel_path)
