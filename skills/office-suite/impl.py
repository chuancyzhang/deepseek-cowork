import os
import json
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader
from core.env_utils import ensure_package_installed

# Lazy import helpers
def get_openpyxl():
    ensure_package_installed("openpyxl")
    import openpyxl
    return openpyxl

def _get_safe_path(workspace_dir, path):
    if not workspace_dir:
        raise ValueError("Workspace not selected.")
    
    abs_path = os.path.abspath(os.path.join(workspace_dir, path))
    abs_workspace = os.path.abspath(workspace_dir)
    
    if not abs_path.startswith(abs_workspace):
        raise ValueError("Access denied (Path Traversal).")
    
    return abs_path

def read_docx(workspace_dir, path):
    """
    Read text content from a DOCX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the DOCX file.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        if not os.path.exists(abs_path):
            return f"Error: File '{path}' does not exist."
            
        doc = Document(abs_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def write_docx(workspace_dir, path, content, mode='w'):
    """
    Write content to a DOCX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the DOCX file.
        content (str): Text content to write.
        mode (str): 'w' to overwrite/create, 'a' to append.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        
        if mode == 'a' and os.path.exists(abs_path):
            doc = Document(abs_path)
        else:
            doc = Document()
            
        # Split by newlines and add as paragraphs
        for line in content.split('\n'):
            doc.add_paragraph(line)
            
        doc.save(abs_path)
        return f"Success: Written to '{path}'."
    except Exception as e:
        return f"Error writing DOCX: {str(e)}"

def read_pptx(workspace_dir, path):
    """
    Read text content from a PPTX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the PPTX file.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        if not os.path.exists(abs_path):
            return f"Error: File '{path}' does not exist."
            
        prs = Presentation(abs_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
            
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def create_pptx(workspace_dir, path, slides_data):
    """
    Create a PPTX file with given slides.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the output PPTX file.
        slides_data (list): List of dicts, each with 'title' and 'content'.
                            Example: [{"title": "Slide 1", "content": "Hello World"}]
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        prs = Presentation()
        
        # Ensure slides_data is a list
        if isinstance(slides_data, str):
            try:
                slides_data = json.loads(slides_data)
            except:
                return "Error: slides_data must be a JSON list or valid list object."

        for slide_info in slides_data:
            title_text = slide_info.get('title', '')
            content_text = slide_info.get('content', '')
            
            # Use a standard layout (Title and Content)
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = title_text
            content.text = content_text
            
        prs.save(abs_path)
        return f"Success: Created presentation at '{path}'."
    except Exception as e:
        return f"Error creating PPTX: {str(e)}"

def read_excel(workspace_dir, path, sheet_name=None):
    """
    Read data from an Excel file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the XLSX file.
        sheet_name (str): Optional sheet name to read.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        if not os.path.exists(abs_path):
            return f"Error: File '{path}' does not exist."
            
        # Use openpyxl for lightweight reading
        openpyxl = get_openpyxl()
        wb = openpyxl.load_workbook(abs_path, data_only=True)
        
        if sheet_name:
            if sheet_name not in wb.sheetnames:
                 return f"Error: Sheet '{sheet_name}' not found. Available: {wb.sheetnames}"
            sheet = wb[sheet_name]
        else:
            sheet = wb.active
            
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # Convert None to empty string for better display
            cleaned_row = [str(cell) if cell is not None else "" for cell in row]
            rows.append("\t".join(cleaned_row))
            
        return "\n".join(rows)
    except Exception as e:
        return f"Error reading Excel: {str(e)}"

def write_excel(workspace_dir, path, data, sheet_name='Sheet1'):
    """
    Write data to an Excel file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the XLSX file.
        data (list): List of lists representing rows.
        sheet_name (str): Name of the sheet.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        
        # Ensure data is a list
        if isinstance(data, str):
            try:
                data = json.loads(data)
            except:
                return "Error: data must be a JSON list of lists."
        
        openpyxl = get_openpyxl()
        
        # Check if file exists to append or create
        if os.path.exists(abs_path):
             wb = openpyxl.load_workbook(abs_path)
             if sheet_name in wb.sheetnames:
                 # If sheet exists, maybe we should clear it or append?
                 # For simplicity, let's create a new sheet if it exists or overwrite?
                 # Let's remove the old sheet and create new one to match 'overwrite' behavior
                 del wb[sheet_name]
             ws = wb.create_sheet(sheet_name)
        else:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
        
        for row in data:
            ws.append(row)
            
        wb.save(abs_path)
        return f"Success: Written to '{path}'."
    except Exception as e:
        return f"Error writing Excel: {str(e)}"

def read_pdf(workspace_dir, path):
    """
    Read text from a PDF file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the PDF file.
    """
    try:
        abs_path = _get_safe_path(workspace_dir, path)
        if not os.path.exists(abs_path):
            return f"Error: File '{path}' does not exist."
            
        reader = PdfReader(abs_path)
        text_content = []
        
        for i, page in enumerate(reader.pages):
            text_content.append(f"--- Page {i+1} ---\n" + page.extract_text())
            
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"
