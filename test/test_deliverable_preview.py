import os
import tempfile
import unittest
from urllib.parse import quote

from core.deliverable_preview import (
    DELIVERABLE_TYPES,
    OFFICE_EXTENSIONS,
    iter_workspace_file_paths,
    linkify_workspace_paths_in_html,
    normalize_workspace_file,
    render_pdf_text_preview,
    render_structured_document_preview,
)


class TestDeliverablePreviewHelpers(unittest.TestCase):
    def _markdown_href(self, path, encode=False):
        href = os.path.normpath(path).replace("\\", "/")
        if len(href) >= 3 and href[1:3] == ":/":
            href = "/" + href
        if encode:
            href = quote(href, safe="/:")
        return href

    def test_supports_modern_and_legacy_office_formats(self):
        for extension in (".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"):
            self.assertIn(extension, DELIVERABLE_TYPES)
            self.assertIn(extension, OFFICE_EXTENSIONS)
        self.assertIn(".md", DELIVERABLE_TYPES)
        self.assertIn(".markdown", DELIVERABLE_TYPES)
        self.assertIn(".bmp", DELIVERABLE_TYPES)
        for extension in (".txt", ".json", ".xml", ".yaml", ".yml", ".log", ".csv", ".tsv"):
            self.assertIn(extension, DELIVERABLE_TYPES)

    def test_only_accepts_existing_workspace_files(self):
        with tempfile.TemporaryDirectory() as workspace, tempfile.TemporaryDirectory() as outside:
            inside_path = os.path.join(workspace, "交付 报告.pdf")
            outside_path = os.path.join(outside, "outside.pdf")
            with open(inside_path, "wb") as handle:
                handle.write(b"pdf")
            with open(outside_path, "wb") as handle:
                handle.write(b"pdf")

            self.assertEqual(normalize_workspace_file(inside_path, workspace), os.path.normpath(inside_path))
            self.assertEqual(normalize_workspace_file(outside_path, workspace), "")
            self.assertEqual(normalize_workspace_file(os.path.join(workspace, "missing.pdf"), workspace), "")

    def test_linkifies_plain_and_inline_code_paths_but_not_outside_files(self):
        with tempfile.TemporaryDirectory() as workspace, tempfile.TemporaryDirectory() as outside:
            report = os.path.join(workspace, "交付 报告.docx")
            external = os.path.join(outside, "external.docx")
            for path in (report, external):
                with open(path, "wb") as handle:
                    handle.write(b"office")
            source = f"<p>完成：{report}</p><code>{report}</code><p>{external}</p>"

            rendered = linkify_workspace_paths_in_html(source, workspace)

            self.assertEqual(rendered.count("cowork-file:"), 2)
            self.assertIn("<code><a ", rendered)
            self.assertIn(external, rendered)

    def test_finds_multiple_supported_paths(self):
        with tempfile.TemporaryDirectory() as workspace:
            paths = [os.path.join(workspace, "report.md"), os.path.join(workspace, "slides.ppt"), os.path.join(workspace, "chart.bmp")]
            for path in paths:
                with open(path, "wb") as handle:
                    handle.write(b"x")
            matches = iter_workspace_file_paths("文件：" + "；".join(paths), workspace)
            self.assertEqual([item[2] for item in matches], [os.path.normpath(path) for path in paths])

    def test_accepts_inline_code_path_but_ignores_fenced_code_path(self):
        with tempfile.TemporaryDirectory() as workspace:
            path = os.path.join(workspace, "report.html")
            with open(path, "wb") as handle:
                handle.write(b"x")

            text = f"行内：`{path}`\n\n```\n{path}\n```"
            matches = iter_workspace_file_paths(text, workspace)

            self.assertEqual([item[2] for item in matches], [os.path.normpath(path)])

    def test_linkifies_complete_path_inside_inline_code(self):
        with tempfile.TemporaryDirectory() as workspace:
            path = os.path.join(workspace, "html_test_output_20260625_225209.pptx")
            with open(path, "wb") as handle:
                handle.write(b"x")

            rendered = linkify_workspace_paths_in_html(f"<p>主文件路径：<code>{path}</code></p>", workspace)

            self.assertIn("cowork-file:", rendered)
            self.assertIn('data-cowork-path=', rendered)

    def test_finds_markdown_link_target_with_leading_windows_slash(self):
        with tempfile.TemporaryDirectory() as workspace:
            path = os.path.join(workspace, "deck.pptx")
            with open(path, "wb") as handle:
                handle.write(b"x")

            text = f"[deck](<{self._markdown_href(path)}>)"
            matches = iter_workspace_file_paths(text, workspace)

            self.assertEqual([item[2] for item in matches], [os.path.normpath(path)])

    def test_rewrites_existing_markdown_file_anchor_to_preview_link(self):
        with tempfile.TemporaryDirectory() as workspace:
            path = os.path.join(workspace, "deck.pptx")
            with open(path, "wb") as handle:
                handle.write(b"x")

            rendered = linkify_workspace_paths_in_html(
                f'<p><a href="{self._markdown_href(path)}">deck</a></p>',
                workspace,
            )

            self.assertIn('href="cowork-file:', rendered)
            self.assertIn(">deck</a>", rendered)
            self.assertIn('data-cowork-path=', rendered)

    def test_rewrites_url_encoded_anchor_with_spaces_and_cjk(self):
        with tempfile.TemporaryDirectory() as workspace:
            path = os.path.join(workspace, "交付 报告.pptx")
            with open(path, "wb") as handle:
                handle.write(b"x")

            rendered = linkify_workspace_paths_in_html(
                f'<p><a href="{self._markdown_href(path, encode=True)}">交付报告</a></p>',
                workspace,
            )

            self.assertIn("cowork-file:", rendered)
            self.assertIn("交付报告", rendered)

    def test_does_not_rewrite_external_missing_or_unsupported_anchors(self):
        with tempfile.TemporaryDirectory() as workspace, tempfile.TemporaryDirectory() as outside:
            outside_path = os.path.join(outside, "outside.pptx")
            unsupported_path = os.path.join(workspace, "notes.bin")
            missing_path = os.path.join(workspace, "missing.pptx")
            for path in (outside_path, unsupported_path):
                with open(path, "wb") as handle:
                    handle.write(b"x")
            source = (
                '<p><a href="https://example.com/report.pptx">web</a>'
                f'<a href="{self._markdown_href(outside_path)}">outside</a>'
                f'<a href="{self._markdown_href(unsupported_path)}">bin</a>'
                f'<a href="{self._markdown_href(missing_path)}">missing</a></p>'
            )

            rendered = linkify_workspace_paths_in_html(source, workspace)

            self.assertNotIn("cowork-file:", rendered)
            self.assertIn("https://example.com/report.pptx", rendered)

    def test_renders_docx_preview_without_microsoft_office(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as workspace:
            source = os.path.join(workspace, "report.docx")
            document = Document()
            document.add_paragraph("季度总结")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "指标"
            table.cell(0, 1).text = "结果"
            document.save(source)

            preview = render_structured_document_preview(source)

        self.assertEqual(preview["format"], "DOCX")
        self.assertIn("季度总结", preview["html"])
        self.assertIn("指标", preview["text"])

    def test_renders_pptx_preview_without_microsoft_office(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as workspace:
            source = os.path.join(workspace, "deck.pptx")
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = "方案页"
            slide.placeholders[1].text = "关键动作"
            presentation.save(source)

            preview = render_structured_document_preview(source)

        self.assertEqual(preview["format"], "PPTX")
        self.assertIn("方案页", preview["html"])
        self.assertIn("关键动作", preview["text"])

    def test_renders_xlsx_preview_without_microsoft_office(self):
        import openpyxl

        with tempfile.TemporaryDirectory() as workspace:
            source = os.path.join(workspace, "data.xlsx")
            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet.title = "明细"
            worksheet.append(["姓名", "金额"])
            worksheet.append(["Alice", 42])
            workbook.save(source)

            preview = render_structured_document_preview(source)

        self.assertEqual(preview["format"], "XLSX")
        self.assertIn("明细", preview["html"])
        self.assertIn("Alice", preview["text"])

    def test_pdf_text_preview_works_without_qtpdf(self):
        from pypdf import PdfWriter

        with tempfile.TemporaryDirectory() as workspace:
            source = os.path.join(workspace, "blank.pdf")
            writer = PdfWriter()
            writer.add_blank_page(width=72, height=72)
            with open(source, "wb") as handle:
                writer.write(handle)

            preview = render_pdf_text_preview(source)

        self.assertIn("PDF 文本预览", preview)
        self.assertIn("Page 1", preview)

    def test_legacy_binary_office_formats_report_clear_error(self):
        with tempfile.TemporaryDirectory() as workspace:
            source = os.path.join(workspace, "old.doc")
            with open(source, "wb") as handle:
                handle.write(b"legacy")

            with self.assertRaisesRegex(RuntimeError, "旧版二进制 Office"):
                render_structured_document_preview(source)


if __name__ == "__main__":
    unittest.main()
