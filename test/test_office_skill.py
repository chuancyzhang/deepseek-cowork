import importlib.util
import json
import os
import shutil
import sys
import tempfile
import unittest

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

file_spec = importlib.util.spec_from_file_location(
    "file_impl",
    os.path.join(os.path.dirname(__file__), "../skills/file-system/impl.py"),
)
file_impl = importlib.util.module_from_spec(file_spec)
file_spec.loader.exec_module(file_impl)

doc_spec = importlib.util.spec_from_file_location(
    "document_impl",
    os.path.join(os.path.dirname(__file__), "../ai_skills/document-reader/impl.py"),
)
document_impl = importlib.util.module_from_spec(doc_spec)
doc_spec.loader.exec_module(document_impl)


class TestOfficeSkill(unittest.TestCase):
    def setUp(self):
        self.workspace_dir = tempfile.mkdtemp()

    def tearDown(self):
        shutil.rmtree(self.workspace_dir, ignore_errors=True)

    def test_plain_text_tool_refuses_structured_documents(self):
        result = file_impl.text_file_read(self.workspace_dir, "test.docx")
        self.assertEqual(result["error"]["code"], "structured_document_not_supported")

        result = file_impl.apply_patch(
            self.workspace_dir,
            "*** Begin Patch\n*** Add File: test.pdf\n+content\n*** End Patch",
        )
        self.assertEqual(result["error"]["code"], "structured_document_not_supported")

    def test_document_read_handles_docx_xlsx_and_plain_text_stays_separate(self):
        from docx import Document
        import openpyxl

        doc = Document()
        doc.add_paragraph("Hello World")
        doc.save(os.path.join(self.workspace_dir, "test.docx"))

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["Name", "Age"])
        sheet.append(["Alice", 30])
        workbook.save(os.path.join(self.workspace_dir, "test.xlsx"))

        with open(os.path.join(self.workspace_dir, "test.txt"), "w", encoding="utf-8") as handle:
            handle.write("Just some text")

        docx_content = document_impl.document_read(self.workspace_dir, "test.docx")
        self.assertIn("Hello World", docx_content["content"])

        xlsx_content = document_impl.document_read(self.workspace_dir, "test.xlsx")
        self.assertIn("Alice", xlsx_content["content"])

        text_content = file_impl.text_file_read(self.workspace_dir, "test.txt")
        self.assertEqual(text_content["content"], "Just some text")

    def test_document_reader_is_default_on_without_runtime_dependency_install(self):
        manifest_path = os.path.join(
            os.path.dirname(__file__),
            "../ai_skills/document-reader/skill.json",
        )
        with open(manifest_path, "r", encoding="utf-8") as handle:
            manifest = json.load(handle)
        with open(document_impl.__file__, "r", encoding="utf-8") as handle:
            implementation = handle.read()

        self.assertTrue(manifest["default_enabled"])
        self.assertEqual(manifest["python_dependencies"], [])
        self.assertNotIn("ensure_package_installed", implementation)

    def test_explicit_disable_still_overrides_document_reader_default(self):
        from core.config_manager import ConfigManager

        config = object.__new__(ConfigManager)
        config.config = {
            "disabled_skills": ["document-reader"],
            "enabled_skills": [],
        }
        self.assertFalse(config.is_skill_enabled("document-reader", default_enabled=True))

    def test_document_read_handles_pptx_and_pdf(self):
        from pptx import Presentation
        from pypdf import PdfWriter

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Quarterly Review"
        presentation.save(os.path.join(self.workspace_dir, "test.pptx"))

        writer = PdfWriter()
        writer.add_blank_page(width=200, height=200)
        with open(os.path.join(self.workspace_dir, "test.pdf"), "wb") as handle:
            writer.write(handle)

        pptx_content = document_impl.document_read(self.workspace_dir, "test.pptx")
        self.assertIn("Quarterly Review", pptx_content["content"])
        self.assertEqual(pptx_content["metadata"]["slide_count"], 1)

        pdf_content = document_impl.document_read(self.workspace_dir, "test.pdf")
        self.assertEqual(pdf_content["format"], "pdf")
        self.assertEqual(pdf_content["metadata"]["total_pages"], 1)


if __name__ == "__main__":
    unittest.main()
